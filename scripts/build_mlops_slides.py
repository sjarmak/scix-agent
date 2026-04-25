"""Build the MLOps Community 'Making Scientific Knowledge Navigable for Agents' deck.

Writes docs/slides/mlops_community_research_copilot.pptx. Run from the repo root:

    python scripts/build_mlops_slides.py

Styling mirrors the earlier deck: Inter, white background, dark ink (#151B23)
for headings, gray (#4B5563) for body, blue (#0B4A9E / #1F6FEB) for accents.
"""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Sequence

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Pt

# Canvas: 13.333" x 7.5" widescreen (same as the previous deck).
SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)

MARGIN_X = Emu(457200)          # 0.5"
FOOTER_Y = Emu(6492240)

FONT = "Inter"

INK = RGBColor(0x15, 0x1B, 0x23)
BODY = RGBColor(0x4B, 0x55, 0x63)
MUTED = RGBColor(0x6B, 0x72, 0x80)
ACCENT = RGBColor(0x0B, 0x4A, 0x9E)
ACCENT_BRIGHT = RGBColor(0x1F, 0x6F, 0xEB)
DIVIDER = RGBColor(0xE5, 0xE7, 0xEB)


@dataclass(frozen=True)
class Deck:
    title: str = "Making Scientific Knowledge Navigable for Agents"
    subtitle_short: str = "Making Scientific Knowledge Navigable for Agents"
    speaker: str = "Stephanie Jarmak"
    venue: str = "MLOps Community · Research Copilot Workshop"
    affiliation: str = "Sourcegraph  ·  Harvard & Smithsonian | Center for Astrophysics"


DECK = Deck()

ASSETS = Path(__file__).resolve().parent.parent / "docs" / "slides_assets"


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _blank_slide(prs: Presentation):
    # Layout 6 is "Blank" in the default master.
    return prs.slides.add_slide(prs.slide_layouts[6])


def _add_text(
    slide,
    *,
    left,
    top,
    width,
    height,
    text,
    font_size,
    bold=False,
    color=INK,
    align="left",
    line_spacing: float | None = None,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)

    lines = text.split("\n") if isinstance(text, str) else list(text)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if line_spacing is not None:
            p.line_spacing = line_spacing
        if align == "center":
            from pptx.enum.text import PP_ALIGN
            p.alignment = PP_ALIGN.CENTER
        elif align == "right":
            from pptx.enum.text import PP_ALIGN
            p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = line
        run.font.name = FONT
        run.font.size = font_size
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def _add_rect(slide, *, left, top, width, height, fill=None, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Emu(6350)
    return shape


def _fit_picture(slide, path: Path, *, box_left, box_top, box_w, box_h, align="center"):
    """Add image scaled to fit inside (box_w, box_h) while preserving aspect.

    Centers inside the box so text placed immediately below `box_top + box_h`
    can never overlap the image.
    """
    from PIL import Image
    with Image.open(path) as im:
        iw, ih = im.size
    aspect = ih / iw
    # Scale to whichever dimension is binding.
    target_w = box_w
    target_h = int(target_w * aspect)
    if target_h > box_h:
        target_h = box_h
        target_w = int(target_h / aspect)
    if align == "center":
        left = box_left + (box_w - target_w) // 2
    elif align == "left":
        left = box_left
    else:
        left = box_left + (box_w - target_w)
    top = box_top + (box_h - target_h) // 2
    slide.shapes.add_picture(str(path), left, top, width=target_w, height=target_h)


def _footer(slide, *, page_number: int, total_pages: int):
    _add_text(
        slide,
        left=MARGIN_X,
        top=FOOTER_Y,
        width=Emu(9144000),
        height=Emu(274320),
        text=f"{DECK.speaker}  ·  MLOps Community  ·  {DECK.subtitle_short}",
        font_size=Pt(10),
        color=MUTED,
    )
    _add_text(
        slide,
        left=Emu(11247120),
        top=FOOTER_Y,
        width=Emu(731520),
        height=Emu(274320),
        text=f"{page_number} / {total_pages}",
        font_size=Pt(10),
        color=MUTED,
        align="right",
    )


def _title_block(slide, *, title: str, subtitle: str | None = None):
    _add_text(
        slide,
        left=MARGIN_X,
        top=Emu(502920),
        width=Emu(11247120),
        height=Emu(731520),
        text=title,
        font_size=Pt(32),
        bold=True,
        color=INK,
    )
    if subtitle:
        _add_text(
            slide,
            left=MARGIN_X,
            top=Emu(1371600),
            width=Emu(11247120),
            height=Emu(457200),
            text=subtitle,
            font_size=Pt(16),
            color=BODY,
        )


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------


def slide_title(prs):
    slide = _blank_slide(prs)

    _add_rect(
        slide,
        left=MARGIN_X,
        top=Emu(1680000),
        width=Emu(1371600),
        height=Emu(76200),
        fill=ACCENT,
    )

    _add_text(
        slide,
        left=Emu(822960),
        top=Emu(1828800),
        width=Emu(10972800),
        height=Emu(1097280),
        text="Making Scientific Knowledge",
        font_size=Pt(52),
        bold=True,
        color=INK,
    )
    _add_text(
        slide,
        left=Emu(822960),
        top=Emu(2926080),
        width=Emu(10972800),
        height=Emu(1097280),
        text="Navigable for Agents",
        font_size=Pt(52),
        bold=True,
        color=INK,
    )
    _add_text(
        slide,
        left=Emu(822960),
        top=Emu(4572000),
        width=Emu(10972800),
        height=Emu(457200),
        text=DECK.speaker,
        font_size=Pt(22),
        bold=True,
        color=INK,
    )
    _add_text(
        slide,
        left=Emu(822960),
        top=Emu(4983480),
        width=Emu(10972800),
        height=Emu(411480),
        text=DECK.affiliation,
        font_size=Pt(16),
        color=BODY,
    )
    _add_text(
        slide,
        left=Emu(822960),
        top=Emu(6035040),
        width=Emu(10972800),
        height=Emu(365760),
        text=DECK.venue,
        font_size=Pt(14),
        color=ACCENT,
    )
    return slide


def slide_reality_of_science(prs, page, total):
    slide = _blank_slide(prs)
    _title_block(
        slide,
        title="Scientific knowledge is fragmented and weakly linked",
        subtitle="We query it as if it were flat text. We can do better.",
    )

    # Four bullets across the width
    bullets = [
        (
            "Knowledge is scattered",
            "Literature, datasets, software, posters, talks — in different systems, "
            "with different identifiers, and no common seams.",
        ),
        (
            "The structure we have is incomplete",
            "DOIs, citation graphs, data-availability statements, ORCIDs — partial "
            "coverage, inconsistent metadata, ambiguous attribution.",
        ),
        (
            "Relationships are underspecified",
            "A citation can mean: builds on, disproves, mentions, reuses data from, "
            "or reuses methods from. The graph collapses all five into one edge.",
        ),
        (
            "Search is the universal API",
            "Keyword or vector search returns a ranked list. That's fine for 'find me "
            "a paper.' It's not enough for 'map the landscape' or 'what breaks if "
            "this result is wrong.'",
        ),
    ]

    top = Emu(2100000)
    row_h = Emu(1000000)
    for i, (head, body) in enumerate(bullets):
        y = top + i * row_h
        _add_rect(
            slide,
            left=MARGIN_X,
            top=y + Emu(60000),
            width=Emu(60000),
            height=Emu(600000),
            fill=ACCENT,
        )
        _add_text(
            slide,
            left=Emu(640000),
            top=y,
            width=Emu(4500000),
            height=Emu(365760),
            text=head,
            font_size=Pt(18),
            bold=True,
            color=INK,
        )
        _add_text(
            slide,
            left=Emu(640000),
            top=y + Emu(380000),
            width=Emu(11000000),
            height=Emu(600000),
            text=body,
            font_size=Pt(14),
            color=BODY,
        )

    _footer(slide, page_number=page, total_pages=total)


def slide_scale_problem(prs, page, total):
    slide = _blank_slide(prs)
    _title_block(
        slide,
        title="The scale problem: research output accelerates like codebases do",
        subtitle="Engineers gave their agents structure. Scientists haven't — yet.",
    )

    stats = [
        ("3M+", "scientific papers published per year", "doubling every ~12 years"),
        ("20–30%", "of a researcher's time", "spent on literature search"),
        ("AI-accelerated", "both code and research output",
         "the same growth curve hits both"),
    ]
    stat_y = Emu(2011680)
    col_w = Emu(3500000)
    gap = Emu(350000)
    for i, (big, mid, small) in enumerate(stats):
        x = MARGIN_X + i * (col_w + gap)
        _add_text(
            slide,
            left=x,
            top=stat_y,
            width=col_w,
            height=Emu(1188720),
            text=big,
            font_size=Pt(56),
            bold=True,
            color=ACCENT,
        )
        _add_text(
            slide,
            left=x,
            top=stat_y + Emu(1188720),
            width=col_w,
            height=Emu(457200),
            text=mid,
            font_size=Pt(16),
            color=INK,
        )
        _add_text(
            slide,
            left=x,
            top=stat_y + Emu(1645920),
            width=col_w,
            height=Emu(365760),
            text=small,
            font_size=Pt(12),
            color=MUTED,
        )

    _add_text(
        slide,
        left=MARGIN_X,
        top=Emu(5050000),
        width=Emu(11247120),
        height=Emu(1200000),
        text=(
            "Codebases outran human reading years ago. Engineers responded by giving "
            "agents structure — dependency graphs, call graphs, blast-radius tools. "
            "Scientific literature has the same shape of problem and has not yet "
            "had the same response."
        ),
        font_size=Pt(16),
        color=BODY,
        line_spacing=1.25,
    )
    _footer(slide, page_number=page, total_pages=total)


def slide_fragmentation(prs, page, total):
    slide = _blank_slide(prs)
    _title_block(
        slide,
        title="Fragmentation of knowledge: serendipity over systems",
        subtitle="Right now, transferable methods travel by chance. That's the thing to fix.",
    )

    left_col = [
        ("Citations encode more than dependency",
         "Prestige, convention, self-citation, availability bias. A clean "
         "technical DAG is not what the citation graph actually is."),
        ("Disciplinary silos hide transferable methods",
         "A technique useful in one field may have been solved in a "
         "neighbouring one years earlier — under different vocabulary."),
    ]
    right_col = [
        ("Connections happen through serendipity, not systems",
         "Whether two lines of work meet depends on who read what, when. "
         "It is not a structural property of the knowledge itself."),
        ("The shift worth making",
         "Turn serendipitous discovery into retrievable relevance — so an "
         "agent can surface it, and a researcher can rely on it."),
    ]

    y0 = Emu(2150000)
    row_h = Emu(1900000)
    cols = [(MARGIN_X, left_col), (Emu(6400000), right_col)]
    for x, col in cols:
        for i, (head, body) in enumerate(col):
            y = y0 + i * row_h
            _add_text(
                slide,
                left=x,
                top=y,
                width=Emu(5400000),
                height=Emu(400000),
                text=head,
                font_size=Pt(17),
                bold=True,
                color=INK,
            )
            _add_text(
                slide,
                left=x,
                top=y + Emu(450000),
                width=Emu(5400000),
                height=Emu(1400000),
                text=body,
                font_size=Pt(14),
                color=BODY,
                line_spacing=1.25,
            )

    _footer(slide, page_number=page, total_pages=total)


def slide_codebase_analogy(prs, page, total):
    slide = _blank_slide(prs)
    _title_block(
        slide,
        title="Codebases as analogy",
        subtitle="Software stopped treating code as isolated files. Science can do the same.",
    )

    # Three parallel rows: software -> science mapping
    rows = [
        ("Dependencies", "what code relies on to run",
         "Reused methods and datasets", "what a result relies on to hold"),
        ("Call graphs", "chains of function calls",
         "Chains of reasoning", "a claim that rests on a claim that rests on a dataset"),
        ("Blast radius", "what breaks if this changes",
         "Result impact", "which follow-on findings are at risk if this one is wrong"),
    ]

    header_y = Emu(1900000)
    _add_text(
        slide, left=MARGIN_X, top=header_y, width=Emu(5400000), height=Emu(320000),
        text="Software engineering", font_size=Pt(14), bold=True, color=MUTED,
    )
    _add_text(
        slide, left=Emu(6400000), top=header_y, width=Emu(5400000), height=Emu(320000),
        text="Science", font_size=Pt(14), bold=True, color=ACCENT,
    )

    row_h = Emu(1150000)
    y0 = Emu(2350000)
    for i, (sw_head, sw_body, sci_head, sci_body) in enumerate(rows):
        y = y0 + i * row_h
        _add_rect(
            slide,
            left=MARGIN_X,
            top=y - Emu(30000),
            width=Emu(11247120),
            height=Emu(6350),
            fill=DIVIDER,
        )
        _add_text(
            slide, left=MARGIN_X, top=y, width=Emu(5400000), height=Emu(360000),
            text=sw_head, font_size=Pt(17), bold=True, color=INK,
        )
        _add_text(
            slide, left=MARGIN_X, top=y + Emu(400000), width=Emu(5400000),
            height=Emu(600000), text=sw_body, font_size=Pt(13), color=BODY,
        )
        _add_text(
            slide, left=Emu(6400000), top=y, width=Emu(5400000), height=Emu(360000),
            text=sci_head, font_size=Pt(17), bold=True, color=INK,
        )
        _add_text(
            slide, left=Emu(6400000), top=y + Emu(400000), width=Emu(5400000),
            height=Emu(600000), text=sci_body, font_size=Pt(13), color=BODY,
        )

    _add_text(
        slide,
        left=MARGIN_X,
        top=Emu(5900000),
        width=Emu(11247120),
        height=Emu(500000),
        text=(
            "Citation graphs approximate all of this — but they're incomplete, "
            "semantically weak, and socially biased. The real signal is the "
            "underlying ecosystem of data, methods, results, code, and claims."
        ),
        font_size=Pt(14),
        color=BODY,
        line_spacing=1.2,
    )
    _footer(slide, page_number=page, total_pages=total)


def slide_debugging_science(prs, page, total):
    slide = _blank_slide(prs)
    _title_block(
        slide,
        title="Toward 'debugging' science",
        subtitle="Treat outputs as claims to evaluate — independent of citation count or venue prestige.",
    )

    top_box = [
        (
            "Science already has implicit acceptance criteria",
            [
                "Physical consistency — does it fit known laws?",
                "Reproducibility — can someone else hit the same result?",
                "Compatibility — does it agree with adjacent findings?",
            ],
        ),
    ]

    y = Emu(1950000)
    head, items = top_box[0]
    _add_text(
        slide, left=MARGIN_X, top=y, width=Emu(11247120), height=Emu(400000),
        text=head, font_size=Pt(18), bold=True, color=INK,
    )
    for i, line in enumerate(items):
        _add_text(
            slide,
            left=MARGIN_X + Emu(200000),
            top=y + Emu(500000) + i * Emu(380000),
            width=Emu(11000000),
            height=Emu(365760),
            text="·  " + line,
            font_size=Pt(14),
            color=BODY,
        )

    # What an agent-navigable layer unlocks
    unlock_y = Emu(3900000)
    _add_text(
        slide, left=MARGIN_X, top=unlock_y, width=Emu(11247120), height=Emu(400000),
        text="What a navigable layer unlocks",
        font_size=Pt(18), bold=True, color=ACCENT,
    )

    unlocks = [
        ("Automated consistency checks",
         "Flag claims that contradict well-established adjacent results."),
        ("Conflicting-result detection",
         "Surface studies that disagree rather than burying the tension in citations."),
        ("Missing-link identification",
         "Point out papers that should cite each other but don't — and why."),
        ("Method-reuse surfacing",
         "Find transferable techniques across silos, not just nearest-neighbor prose."),
    ]
    col_w = Emu(5400000)
    ry = unlock_y + Emu(500000)
    for i, (h, b) in enumerate(unlocks):
        col = i % 2
        row = i // 2
        x = MARGIN_X + col * (col_w + Emu(450000))
        yy = ry + row * Emu(950000)
        _add_text(
            slide, left=x, top=yy, width=col_w, height=Emu(340000),
            text=h, font_size=Pt(14), bold=True, color=INK,
        )
        _add_text(
            slide, left=x, top=yy + Emu(360000), width=col_w, height=Emu(600000),
            text=b, font_size=Pt(12), color=BODY, line_spacing=1.2,
        )

    _footer(slide, page_number=page, total_pages=total)


def slide_navigable_means(prs, page, total):
    slide = _blank_slide(prs)
    _title_block(
        slide,
        title='What "navigable" means for an agent',
        subtitle="Agents benefit from ranked lists too — they just benefit from more.",
    )

    left_h = "Human searcher"
    right_h = "Agent (given the right primitives)"

    left_items = [
        ("Ranked list", "Scans titles, opens a few, satisfied — same primitive an agent uses for cheap queries."),
        ("One-shot query", "Session state lives in the researcher's head."),
        ("Visual cues", "Journal, authors, year, formatting."),
        ("Stops when tired", "Implicit budget, intuition-driven."),
    ]
    right_items = [
        ("Graph topology", "Which papers bridge communities? Which are structurally central?"),
        ("Session working set", "Accumulates, dedupes, tags, reasons across turns."),
        ("Structured returns", "Bibcodes, IDs, community labels, provenance — composes with other tools."),
        ("Multi-hop moves", "Citation chains, co-citation, bridge-paper gaps, temporal."),
    ]

    hdr_y = Emu(1828800)
    _add_text(
        slide, left=Emu(548640), top=hdr_y, width=Emu(5303520), height=Emu(457200),
        text=left_h, font_size=Pt(20), bold=True, color=MUTED,
    )
    _add_text(
        slide, left=Emu(6309360), top=hdr_y, width=Emu(5303520), height=Emu(457200),
        text=right_h, font_size=Pt(20), bold=True, color=ACCENT,
    )

    y0 = Emu(2651760)
    row = Emu(860000)
    for i, ((lh, lb), (rh, rb)) in enumerate(zip(left_items, right_items)):
        y = y0 + i * row
        _add_text(
            slide, left=Emu(548640), top=y, width=Emu(5303520), height=Emu(320040),
            text=lh, font_size=Pt(15), bold=True, color=INK,
        )
        _add_text(
            slide, left=Emu(548640), top=y + Emu(320040), width=Emu(5303520),
            height=Emu(502920), text=lb, font_size=Pt(13), color=BODY, line_spacing=1.2,
        )
        _add_text(
            slide, left=Emu(6309360), top=y, width=Emu(5303520), height=Emu(320040),
            text=rh, font_size=Pt(15), bold=True, color=INK,
        )
        _add_text(
            slide, left=Emu(6309360), top=y + Emu(320040), width=Emu(5303520),
            height=Emu(502920), text=rb, font_size=Pt(13), color=BODY, line_spacing=1.2,
        )

    _footer(slide, page_number=page, total_pages=total)


def slide_hybrid_retrieval(prs, page, total):
    slide = _blank_slide(prs)
    _title_block(
        slide,
        title="Vector search and knowledge graphs, together",
        subtitle="Hybrid retrieval is the table stakes. Agent-mediated iteration is the unlock.",
    )

    blocks = [
        ("Semantic", "Dense vectors — domain-tuned plus general-purpose, fused via reciprocal rank."),
        ("Lexical", "BM25 on titles, abstracts, keywords. Catches specific terms that embeddings blur."),
        ("Structural", "Graph moves: citation chains, co-citation, shortest paths, community lenses."),
        ("Symbolic", "Entity and claim extraction — instruments, datasets, methods, results as first-class objects."),
        ("Iterative",
         "An agent chooses the next move based on what it just learned. "
         "The loop — not any single tool — is where the leverage lives."),
    ]

    y = Emu(1950000)
    col_w = Emu(3600000)
    gap = Emu(180000)
    for i, (h, b) in enumerate(blocks[:4]):
        x = MARGIN_X + (i % 3) * (col_w + gap)
        yy = y + (i // 3) * Emu(1850000)
        _add_rect(slide, left=x, top=yy, width=col_w, height=Emu(1600000),
                  fill=RGBColor(0xF8, 0xFA, 0xFC), line=DIVIDER)
        _add_text(slide, left=x + Emu(200000), top=yy + Emu(200000),
                  width=col_w - Emu(400000), height=Emu(400000),
                  text=h, font_size=Pt(18), bold=True, color=ACCENT)
        _add_text(slide, left=x + Emu(200000), top=yy + Emu(650000),
                  width=col_w - Emu(400000), height=Emu(900000),
                  text=b, font_size=Pt(13), color=BODY, line_spacing=1.3)

    # Iterative spans the bottom wide
    iy = y + Emu(1850000)
    x = MARGIN_X + col_w + gap
    _add_rect(slide, left=x, top=iy, width=col_w, height=Emu(1600000),
              fill=ACCENT, line=None)
    _add_text(slide, left=x + Emu(200000), top=iy + Emu(200000),
              width=col_w - Emu(400000), height=Emu(400000),
              text=blocks[4][0], font_size=Pt(18), bold=True,
              color=RGBColor(0xFF, 0xFF, 0xFF))
    _add_text(slide, left=x + Emu(200000), top=iy + Emu(650000),
              width=col_w - Emu(400000), height=Emu(900000),
              text=blocks[4][1], font_size=Pt(13),
              color=RGBColor(0xE5, 0xEE, 0xFB), line_spacing=1.3)

    _footer(slide, page_number=page, total_pages=total)


def slide_why_mcp(prs, page, total):
    slide = _blank_slide(prs)
    _title_block(
        slide,
        title="Why MCP is a good fit for this layer",
        subtitle="Model Context Protocol gives the agent tools it can actually hold.",
    )

    items = [
        ("Tool-calling primitives",
         "First-class: name, typed args, typed returns. The agent knows what it can do."),
        ("Session context",
         "Stateful per-connection. Working sets live on the server, not in the prompt."),
        ("Structured responses",
         "JSON objects, not HTML snippets. Compose cleanly with other tools."),
        ("Transport-agnostic",
         "stdio, HTTP, WebSocket. Works with Claude Desktop, Cursor, custom agents."),
        ("Portable discipline",
         "If your research copilot speaks MCP, the same knowledge layer can be "
         "hosted by ADS, arXiv, your own group, or someone else."),
    ]

    col_w = Emu(5400000)
    y0 = Emu(1950000)
    row_h = Emu(1100000)
    for i, (h, b) in enumerate(items):
        col = i % 2
        row = i // 2
        x = MARGIN_X + col * (col_w + Emu(450000))
        y = y0 + row * row_h
        _add_text(slide, left=x, top=y, width=col_w, height=Emu(360000),
                  text=h, font_size=Pt(16), bold=True, color=INK)
        _add_text(slide, left=x, top=y + Emu(380000), width=col_w, height=Emu(700000),
                  text=b, font_size=Pt(13), color=BODY, line_spacing=1.25)

    _footer(slide, page_number=page, total_pages=total)


def slide_corpus_and_arch(prs, page, total):
    slide = _blank_slide(prs)
    _title_block(
        slide,
        title="A working prototype: SciX on NASA ADS",
        subtitle="32.4M papers, PostgreSQL + pgvector, single server, one MCP endpoint.",
    )

    img_path = ASSETS / "architecture.png"
    # Reserve the bottom 0.65" for the caption so the image can never overlap it.
    CAPTION_TOP = Emu(5900000)
    IMAGE_TOP = Emu(1950000)
    IMAGE_BOTTOM_MAX = CAPTION_TOP - Emu(150000)
    if img_path.exists():
        _fit_picture(
            slide,
            img_path,
            box_left=MARGIN_X,
            box_top=IMAGE_TOP,
            box_w=SLIDE_W - 2 * MARGIN_X,
            box_h=IMAGE_BOTTOM_MAX - IMAGE_TOP,
        )

    _add_text(
        slide,
        left=MARGIN_X,
        top=CAPTION_TOP,
        width=Emu(11247120),
        height=Emu(450000),
        text=(
            "Deliberate simplicity: no separate vector DB, no orchestrator, no "
            "queue. Postgres + pgvector carries the corpus, the graph metrics, "
            "and the retrieval stack."
        ),
        font_size=Pt(13),
        color=BODY,
        line_spacing=1.25,
    )
    _footer(slide, page_number=page, total_pages=total)


def slide_coverage_matters(prs, page, total):
    slide = _blank_slide(prs)
    _title_block(
        slide,
        title="Ask your copilot what it can see",
        subtitle="Corpus boundary is a product decision, not an operational one.",
    )

    points = [
        ("Citation graphs are shape-sensitive",
         "A partial corpus breaks PageRank, community detection, and co-citation — "
         "not by a few percent, but structurally. Recent papers cite older papers; "
         "a rolling window systematically excludes the very papers its contents reference."),
        ("The practical question for a user",
         "What does the copilot's corpus actually include? Where do its edges dangle? "
         "If you ask it 'which papers are foundational here,' is the answer computed "
         "over a 6-year window or over the full field?"),
        ("The engineering answer is boring",
         "Full-corpus ingest is feasible. ADS → Postgres via binary COPY, ~4 hours on "
         "one machine. There is no excuse for partial-graph analytics."),
    ]

    y0 = Emu(1950000)
    row_h = Emu(1350000)
    for i, (h, b) in enumerate(points):
        y = y0 + i * row_h
        _add_text(slide, left=MARGIN_X, top=y, width=Emu(11247120), height=Emu(400000),
                  text=h, font_size=Pt(17), bold=True, color=INK)
        _add_text(slide, left=MARGIN_X, top=y + Emu(430000), width=Emu(11247120),
                  height=Emu(900000), text=b, font_size=Pt(13), color=BODY,
                  line_spacing=1.3)

    _footer(slide, page_number=page, total_pages=total)


def slide_two_vector_stores(prs, page, total):
    slide = _blank_slide(prs)
    _title_block(
        slide,
        title="Two vector stores, two jobs",
        subtitle="Pick the backend for the move — the agent only sees one MCP interface.",
    )

    cols = [
        (
            "pgvector — the corpus",
            [
                "All 32.4M papers",
                "Dense ANN via HNSW, joined with graph metrics and communities",
                "One SQL query composes retrieval + filters + PageRank",
                "Good at: topical search, hybrid RRF, structural filters",
            ],
        ),
        (
            "Qdrant — the recommendation move",
            [
                "Pilot subset (top-N by PageRank) with rich payload",
                "Multi-example recommendation: 'more like these, less like those'",
                "Payload filters: year, doctype, community, arxiv class",
                "Good at: the one move pgvector can't cleanly do in SQL",
            ],
        ),
    ]
    col_w = Emu(5400000)
    gap = Emu(450000)
    y = Emu(1950000)
    for i, (head, bullets) in enumerate(cols):
        x = MARGIN_X + i * (col_w + gap)
        _add_rect(slide, left=x, top=y, width=col_w, height=Emu(3700000),
                  fill=RGBColor(0xF8, 0xFA, 0xFC), line=DIVIDER)
        _add_text(slide, left=x + Emu(220000), top=y + Emu(220000),
                  width=col_w - Emu(440000), height=Emu(440000),
                  text=head, font_size=Pt(18), bold=True, color=ACCENT)
        for j, line in enumerate(bullets):
            _add_text(
                slide,
                left=x + Emu(220000),
                top=y + Emu(900000) + j * Emu(620000),
                width=col_w - Emu(440000),
                height=Emu(580000),
                text="·  " + line,
                font_size=Pt(13),
                color=BODY,
                line_spacing=1.25,
            )

    _add_text(
        slide,
        left=MARGIN_X,
        top=Emu(5900000),
        width=Emu(11247120),
        height=Emu(500000),
        text=(
            "One MCP tool surface (find_similar_by_examples). Two backends. "
            "The agent picks the move by name — not by where the vectors live."
        ),
        font_size=Pt(13),
        color=BODY,
        line_spacing=1.25,
    )
    _footer(slide, page_number=page, total_pages=total)


def slide_tool_design(prs, page, total):
    slide = _blank_slide(prs)
    _title_block(
        slide,
        title="Tool design for agent navigation of scientific literature",
        subtitle="Principles from watching agents actually use these tools.",
    )

    items = [
        ("Fewer, composable tools > many specific tools",
         "Agent tool-selection accuracy degrades past a point. Prefer a small set "
         "of primitives the agent can combine over a catalog of bespoke endpoints."),
        ("Return compact stubs, full text on request",
         "bibcode · title · first_author · year · citation_count · snippet. Full "
         "paper only when the agent asks for it. Context is the scarcest resource."),
        ("Stable, typed schemas",
         "JSON, no free-text fields that need parsing. Every response composes "
         "cleanly with whatever the agent does next."),
        ("Include provenance",
         "Which tool returned this? Which query? The agent can reason about its "
         "own trail — and so can you, when you read the trace."),
        ("Cursor-based pagination",
         "Never return 10,000 bibcodes in one payload. Give the agent a cursor "
         "and let it ask for more if it needs more."),
        ("Idempotent reads, explicit side effects",
         "Same inputs → same outputs. Caching is trivial. Anything that writes "
         "state is called out in the tool name."),
    ]

    col_w = Emu(5400000)
    y0 = Emu(1950000)
    row_h = Emu(1050000)
    for i, (h, b) in enumerate(items):
        col = i % 2
        row = i // 2
        x = MARGIN_X + col * (col_w + Emu(450000))
        y = y0 + row * row_h
        _add_text(slide, left=x, top=y, width=col_w, height=Emu(360000),
                  text=h, font_size=Pt(15), bold=True, color=INK)
        _add_text(slide, left=x, top=y + Emu(380000), width=col_w, height=Emu(700000),
                  text=b, font_size=Pt(13), color=BODY, line_spacing=1.25)

    _footer(slide, page_number=page, total_pages=total)


def slide_unlocks(prs, page, total):
    slide = _blank_slide(prs)
    _title_block(
        slide,
        title="Unlocks: exploration modes you can't do with search alone",
        subtitle="Each of these is a move an agent takes — not a UI the human drives.",
    )

    items = [
        ("Bridge papers between communities",
         "'Which papers are cited by both the exoplanet community and the "
         "atmospheric-chemistry community?' — a single structural query, not a "
         "month of hallway conversations."),
        ("Citation chains (what this claim depends on)",
         "Follow the chain of results the current paper rests on. Where does the "
         "chain terminate — a dataset? a theorem? an assumption someone flagged?"),
        ("Co-citation and bibliographic coupling",
         "Find papers that are read-alike (cited together) or reference-alike "
         "(share references). These are invisible-college detectors."),
        ("Temporal evolution",
         "Which subfields are growing? Which are consolidating into review papers? "
         "Which had a flurry five years ago and went quiet?"),
        ("Multi-lens community assignment",
         "Three lenses — citation, semantic, taxonomic — and their disagreements "
         "are informative. A paper that's semantically central but citationally "
         "isolated is often an 'outsider with a good idea.'"),
    ]

    col_w = Emu(5400000)
    y0 = Emu(1950000)
    row_h = Emu(1050000)
    for i, (h, b) in enumerate(items[:4]):
        col = i % 2
        row = i // 2
        x = MARGIN_X + col * (col_w + Emu(450000))
        y = y0 + row * row_h
        _add_text(slide, left=x, top=y, width=col_w, height=Emu(360000),
                  text=h, font_size=Pt(15), bold=True, color=INK)
        _add_text(slide, left=x, top=y + Emu(380000), width=col_w, height=Emu(700000),
                  text=b, font_size=Pt(12), color=BODY, line_spacing=1.22)

    # The 5th item spans the bottom
    y5 = y0 + 2 * row_h
    h, b = items[4]
    _add_text(slide, left=MARGIN_X, top=y5, width=Emu(11247120), height=Emu(360000),
              text=h, font_size=Pt(15), bold=True, color=INK)
    _add_text(slide, left=MARGIN_X, top=y5 + Emu(380000), width=Emu(11247120),
              height=Emu(500000), text=b, font_size=Pt(12), color=BODY,
              line_spacing=1.22)

    _footer(slide, page_number=page, total_pages=total)


def slide_community_lenses(prs, page, total):
    slide = _blank_slide(prs)
    _title_block(
        slide,
        title="Three community lenses — and their informative disagreements",
        subtitle="Each lens has clean semantics. The agent picks the one that fits the question.",
    )

    cols = [
        (
            "Citation",
            "behavioral affinity",
            [
                "Leiden on the full citation graph",
                "Three resolutions: ~20 · ~200 · ~2000",
                "Coverage: 62% of papers",
                "'Papers that cite each other'",
                "Finds invisible-college patterns",
            ],
        ),
        (
            "Semantic",
            "topical affinity",
            [
                "k-means on INDUS embeddings",
                "Three resolutions",
                "Coverage: 99% (all embedded)",
                "'Papers with similar content'",
                "Catches cross-citing subfields",
            ],
        ),
        (
            "Taxonomic",
            "institutional affinity",
            [
                "arXiv class + UAT concepts",
                "Hierarchical, human-curated",
                "Coverage: ~8% (arXiv subset)",
                "'Papers in the same field'",
                "Sanity-check for the other two",
            ],
        ),
    ]

    col_w = Emu(3700000)
    gap = Emu(200000)
    y = Emu(1950000)
    for i, (name, tag, bullets) in enumerate(cols):
        x = MARGIN_X + i * (col_w + gap)
        _add_rect(slide, left=x, top=y, width=col_w, height=Emu(3800000),
                  fill=RGBColor(0xF8, 0xFA, 0xFC), line=DIVIDER)
        _add_text(slide, left=x + Emu(220000), top=y + Emu(200000),
                  width=col_w - Emu(440000), height=Emu(420000),
                  text=name, font_size=Pt(20), bold=True, color=ACCENT)
        _add_text(slide, left=x + Emu(220000), top=y + Emu(650000),
                  width=col_w - Emu(440000), height=Emu(320000),
                  text=tag, font_size=Pt(12), color=MUTED)
        for j, line in enumerate(bullets):
            _add_text(
                slide,
                left=x + Emu(220000),
                top=y + Emu(1050000) + j * Emu(460000),
                width=col_w - Emu(440000),
                height=Emu(420000),
                text="·  " + line,
                font_size=Pt(12),
                color=BODY,
            )

    _add_text(
        slide,
        left=MARGIN_X,
        top=Emu(5900000),
        width=Emu(11247120),
        height=Emu(500000),
        text=(
            "A paper that's semantically central but citationally isolated is usually "
            "an outsider with a good idea — or a review the field hasn't caught up to."
        ),
        font_size=Pt(13),
        color=BODY,
        line_spacing=1.25,
    )
    _footer(slide, page_number=page, total_pages=total)


def slide_live_demo(prs, page, total):
    slide = _blank_slide(prs)
    _title_block(
        slide,
        title="Live demo",
        subtitle="A research copilot maps the exoplanet-atmospheres landscape.",
    )

    _add_text(
        slide, left=MARGIN_X, top=Emu(1950000), width=Emu(11247120), height=Emu(400000),
        text="User prompt", font_size=Pt(14), bold=True, color=MUTED,
    )
    _add_text(
        slide,
        left=MARGIN_X,
        top=Emu(2300000),
        width=Emu(11247120),
        height=Emu(900000),
        text=(
            "'I'm preparing a grant on exoplanet atmospheric characterization. "
            "Map the research landscape — major subcommunities, foundational papers, "
            "recent bridges between communities, and what's worth reading first.'"
        ),
        font_size=Pt(14),
        color=INK,
        line_spacing=1.3,
    )

    _add_text(
        slide, left=MARGIN_X, top=Emu(3300000), width=Emu(11247120), height=Emu(400000),
        text="Moves the copilot makes", font_size=Pt(14), bold=True, color=MUTED,
    )
    moves = [
        ("search", "broad hybrid hit on 'exoplanet atmospheric characterization'"),
        ("graph_context", "pull community assignments + PageRank for the top hits"),
        ("citation_similarity", "expand with co-citation bridges"),
        ("temporal_evolution", "see which communities are growing"),
        ("find_gaps", "bridge papers between working-set communities"),
        ("read_paper", "deep-dive on 3 foundational candidates"),
    ]
    y = Emu(3700000)
    for i, (tool, desc) in enumerate(moves):
        row_y = y + i * Emu(370000)
        _add_text(
            slide, left=MARGIN_X, top=row_y, width=Emu(350000), height=Emu(320000),
            text=str(i + 1), font_size=Pt(16), bold=True, color=ACCENT_BRIGHT,
        )
        _add_text(
            slide, left=Emu(840000), top=row_y, width=Emu(3000000), height=Emu(320000),
            text=tool, font_size=Pt(14), bold=True, color=INK,
        )
        _add_text(
            slide, left=Emu(4000000), top=row_y, width=Emu(8000000), height=Emu(320000),
            text=desc, font_size=Pt(13), color=BODY,
        )

    _footer(slide, page_number=page, total_pages=total)


def slide_viz(prs, page, total):
    slide = _blank_slide(prs)
    _title_block(
        slide,
        title="What the copilot sees while it works",
        subtitle="Agent-trace overlay · UMAP of 32M INDUS embeddings.",
    )

    # Two images side by side with a caption below them.
    CAPTION_TOP = Emu(5900000)
    IMAGE_TOP = Emu(1950000)
    IMAGE_BOTTOM_MAX = CAPTION_TOP - Emu(150000)
    col_w = Emu(5500000)
    gap = Emu(300000)
    total_w = 2 * col_w + gap
    left_start = (SLIDE_W - total_w) // 2

    trace_path = ASSETS / "agent_trace.png"
    umap_path = ASSETS / "umap_browser.png"
    box_h = IMAGE_BOTTOM_MAX - IMAGE_TOP

    if trace_path.exists():
        _fit_picture(
            slide,
            trace_path,
            box_left=left_start,
            box_top=IMAGE_TOP,
            box_w=col_w,
            box_h=box_h,
        )
    if umap_path.exists():
        _fit_picture(
            slide,
            umap_path,
            box_left=left_start + col_w + gap,
            box_top=IMAGE_TOP,
            box_w=col_w,
            box_h=box_h,
        )

    _add_text(
        slide,
        left=MARGIN_X,
        top=CAPTION_TOP,
        width=Emu(11247120),
        height=Emu(450000),
        text=(
            "These are the structural views an agent reasons over — "
            "not visually, but as graph edges, community IDs, and embedding coordinates."
        ),
        font_size=Pt(13),
        color=BODY,
        line_spacing=1.25,
    )
    _footer(slide, page_number=page, total_pages=total)


def slide_for_users(prs, page, total):
    slide = _blank_slide(prs)
    _title_block(
        slide,
        title="What this means if you're using a research copilot",
        subtitle="Most of you will be integrating one, not building the knowledge layer. "
                 "The questions to ask are the same.",
    )

    items = [
        ("What does the corpus cover?",
         "Full field or a rolling window? Do the graph metrics it returns mean "
         "what you think they mean?"),
        ("What moves does it support beyond search?",
         "Citation chains, co-citation, community bridges, temporal slices. If "
         "the only primitive is search, you have a faster librarian — not a "
         "navigation layer."),
        ("Does it return structured objects?",
         "A bibcode, not a snippet. IDs compose; prose doesn't. Tools that return "
         "prose force the agent to parse its own context."),
        ("Can you trace what it consulted?",
         "Every answer should come with the tool calls it made. This is how you "
         "trust it, debug it, and teach it."),
        ("How does it complement, not replace, your judgement?",
         "Use it to expand the landscape and surface blind spots. Reserve peer "
         "review, interpretation, and prioritization for the human."),
    ]

    col_w = Emu(5400000)
    y0 = Emu(2050000)
    row_h = Emu(1000000)
    for i, (h, b) in enumerate(items[:4]):
        col = i % 2
        row = i // 2
        x = MARGIN_X + col * (col_w + Emu(450000))
        y = y0 + row * row_h
        _add_text(slide, left=x, top=y, width=col_w, height=Emu(360000),
                  text=h, font_size=Pt(15), bold=True, color=INK)
        _add_text(slide, left=x, top=y + Emu(380000), width=col_w, height=Emu(600000),
                  text=b, font_size=Pt(12), color=BODY, line_spacing=1.22)

    y5 = y0 + 2 * row_h
    h, b = items[4]
    _add_text(slide, left=MARGIN_X, top=y5, width=Emu(11247120), height=Emu(360000),
              text=h, font_size=Pt(15), bold=True, color=ACCENT)
    _add_text(slide, left=MARGIN_X, top=y5 + Emu(380000), width=Emu(11247120),
              height=Emu(600000), text=b, font_size=Pt(12), color=BODY,
              line_spacing=1.22)

    _footer(slide, page_number=page, total_pages=total)


def slide_takeaway(prs, page, total):
    slide = _blank_slide(prs)
    _title_block(slide, title="Takeaway")

    lines = [
        ("From search", "to understanding"),
        ("From retrieval", "to navigable scientific context"),
        ("From serendipity", "to systems"),
    ]

    y0 = Emu(2200000)
    row_h = Emu(900000)
    for i, (a, b) in enumerate(lines):
        y = y0 + i * row_h
        _add_text(
            slide, left=MARGIN_X, top=y, width=Emu(4500000), height=Emu(560000),
            text=a, font_size=Pt(30), bold=True, color=MUTED,
        )
        _add_text(
            slide, left=Emu(5200000), top=y, width=Emu(200000), height=Emu(560000),
            text="→", font_size=Pt(30), bold=True, color=MUTED,
        )
        _add_text(
            slide, left=Emu(5550000), top=y, width=Emu(6500000), height=Emu(560000),
            text=b, font_size=Pt(30), bold=True, color=ACCENT,
        )

    _add_text(
        slide,
        left=MARGIN_X,
        top=Emu(5700000),
        width=Emu(11247120),
        height=Emu(500000),
        text=(
            "Scientific knowledge gets more useful when we treat it like the system "
            "it already is — data, methods, results, code, claims — and let agents "
            "navigate the connections we've been leaving to chance."
        ),
        font_size=Pt(14),
        color=BODY,
        line_spacing=1.25,
    )
    _footer(slide, page_number=page, total_pages=total)


def slide_thanks(prs, page, total):
    slide = _blank_slide(prs)
    _add_text(
        slide,
        left=MARGIN_X,
        top=Emu(1800000),
        width=Emu(11247120),
        height=Emu(900000),
        text="Thank you",
        font_size=Pt(52),
        bold=True,
        color=INK,
    )
    _add_text(
        slide,
        left=MARGIN_X,
        top=Emu(2750000),
        width=Emu(11247120),
        height=Emu(500000),
        text="Questions?",
        font_size=Pt(22),
        color=BODY,
    )

    _add_text(
        slide,
        left=MARGIN_X,
        top=Emu(4000000),
        width=Emu(11247120),
        height=Emu(400000),
        text=DECK.speaker,
        font_size=Pt(18),
        bold=True,
        color=INK,
    )
    _add_text(
        slide,
        left=MARGIN_X,
        top=Emu(4400000),
        width=Emu(11247120),
        height=Emu(400000),
        text=DECK.affiliation,
        font_size=Pt(14),
        color=BODY,
    )

    _add_text(
        slide,
        left=MARGIN_X,
        top=Emu(5200000),
        width=Emu(11247120),
        height=Emu(320000),
        text="Resources",
        font_size=Pt(13),
        bold=True,
        color=MUTED,
    )
    resources = [
        "Repository — github.com/sjarmak/scix-agent",
        "Live MCP endpoint — mcp.sjarmak.ai/mcp/  (local-only during workshop)",
        "Paper (ADASS 2026, preprint) — arXiv (pending)",
        "Contact — stephanie.jarmak1@gmail.com",
    ]
    for i, line in enumerate(resources):
        _add_text(
            slide,
            left=MARGIN_X,
            top=Emu(5550000) + i * Emu(280000),
            width=Emu(11247120),
            height=Emu(280000),
            text="·  " + line,
            font_size=Pt(12),
            color=BODY,
        )
    _footer(slide, page_number=page, total_pages=total)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------


def build(output: Path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        slide_title,
        slide_reality_of_science,
        slide_scale_problem,
        slide_fragmentation,
        slide_codebase_analogy,
        slide_debugging_science,
        slide_navigable_means,
        slide_hybrid_retrieval,
        slide_why_mcp,
        slide_corpus_and_arch,
        slide_two_vector_stores,
        slide_tool_design,
        slide_unlocks,
        slide_community_lenses,
        slide_live_demo,
        slide_viz,
        slide_for_users,
        slide_takeaway,
        slide_thanks,
    ]

    total = len(builders)
    for i, builder in enumerate(builders, start=1):
        if i == 1:
            builder(prs)
        else:
            builder(prs, i, total)

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    print(f"wrote {output} ({total} slides)")


if __name__ == "__main__":
    out = Path(__file__).resolve().parent.parent / "docs" / "slides" / "mlops_community_research_copilot.pptx"
    build(out)
