"""Build the MLOps Community slide deck for Stephanie Jarmak.

Talk: "Making Scientific Knowledge Navigable for Agents"
Audience: ML engineers / researchers / data scientists building research copilots
Duration: 30 min (28 slides)
Output: pptx for Google Slides import
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

ASSETS = Path("/home/ds/projects/scix_experiments/docs/slides_assets")
OUT_PATH = Path("/home/ds/projects/scix_experiments/docs/slides/mlops_community_research_copilot.pptx")
OUT_PATH.parent.mkdir(parents=True, exist_ok=True)

# --- design system -----------------------------------------------------------
BG = RGBColor(0xFF, 0xFF, 0xFF)            # white
INK = RGBColor(0x15, 0x1B, 0x23)           # near-black
INK_MUTED = RGBColor(0x4B, 0x55, 0x63)
INK_SOFT = RGBColor(0x6B, 0x72, 0x80)
ACCENT = RGBColor(0x1F, 0x6F, 0xEB)        # blue
ACCENT_DARK = RGBColor(0x0B, 0x4A, 0x9E)
ACCENT_SOFT = RGBColor(0xEA, 0xF2, 0xFE)
GREEN = RGBColor(0x2A, 0xA8, 0x76)
RED = RGBColor(0xD0, 0x56, 0x4A)
RULE = RGBColor(0xE5, 0xE7, 0xEB)

FONT = "Inter"  # Google Slides will fall back gracefully; default feel is clean sans

# 16:9 widescreen
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height


# --- helpers ----------------------------------------------------------------
def add_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_bg(slide, color=BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    return bg


def add_text(slide, text, left, top, width, height, *,
             size=18, bold=False, italic=False, color=INK, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, font=FONT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb


def add_bullets(slide, bullets, left, top, width, height, *,
                size=18, color=INK, bullet_color=ACCENT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        # bullet char as an accent dot
        r = p.add_run()
        r.text = "•  "
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = True
        r.font.color.rgb = bullet_color
        r2 = p.add_run()
        r2.text = b
        r2.font.name = FONT
        r2.font.size = Pt(size)
        r2.font.color.rgb = color
    return tb


def add_rule(slide, left, top, width, height=Pt(2), color=ACCENT):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    r.line.fill.background()
    r.fill.solid()
    r.fill.fore_color.rgb = color
    return r


def add_footer(slide, idx, total):
    add_text(slide, f"Stephanie Jarmak  ·  MLOps Community  ·  Making Scientific Knowledge Navigable for Agents",
             Inches(0.5), Inches(7.1), Inches(10), Inches(0.3),
             size=10, color=INK_SOFT)
    add_text(slide, f"{idx} / {total}",
             Inches(12.3), Inches(7.1), Inches(0.8), Inches(0.3),
             size=10, color=INK_SOFT, align=PP_ALIGN.RIGHT)


def add_title(slide, text, *, sub=None, top=Inches(0.55)):
    add_text(slide, text, Inches(0.5), top, Inches(12.3), Inches(0.8),
             size=32, bold=True, color=INK)
    add_rule(slide, Inches(0.5), Inches(1.35), Inches(0.8))
    if sub:
        add_text(slide, sub, Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.5),
                 size=16, color=INK_MUTED, italic=True)


def add_image(slide, path, left, top, *, width=None, height=None):
    if width is not None and height is not None:
        return slide.shapes.add_picture(str(path), left, top, width=width, height=height)
    if width is not None:
        return slide.shapes.add_picture(str(path), left, top, width=width)
    if height is not None:
        return slide.shapes.add_picture(str(path), left, top, height=height)
    return slide.shapes.add_picture(str(path), left, top)


def set_notes(slide, text):
    nf = slide.notes_slide.notes_text_frame
    nf.text = text


TOTAL_SLIDES = 30

# ============================================================================
#  SLIDE 1 — Title
# ============================================================================
s = add_blank(prs)
set_bg(s)
# left accent bar
add_rule(s, 0, 0, Inches(0.18), height=SH, color=ACCENT)
add_text(s, "Making Scientific Knowledge\nNavigable for Agents",
         Inches(0.9), Inches(2.0), Inches(12), Inches(2.4),
         size=52, bold=True, color=INK)
add_rule(s, Inches(0.9), Inches(4.2), Inches(1.2))
add_text(s, "Stephanie Jarmak",
         Inches(0.9), Inches(4.45), Inches(12), Inches(0.5),
         size=22, bold=True, color=INK)
add_text(s, "Sourcegraph  ·  Harvard & Smithsonian | Center for Astrophysics",
         Inches(0.9), Inches(4.9), Inches(12), Inches(0.45),
         size=16, color=INK_MUTED)
add_text(s, "MLOps Community · Research Copilot Workshop",
         Inches(0.9), Inches(6.6), Inches(12), Inches(0.4),
         size=14, italic=True, color=ACCENT_DARK)
set_notes(s,
    "Welcome. I'm Stephanie Jarmak — I work at Sourcegraph on code-intelligence "
    "infrastructure and I'm a research affiliate at the Harvard-Smithsonian Center "
    "for Astrophysics. Today I want to share what we've learned building an "
    "agent-native layer over the full NASA ADS corpus — 32.4 million scientific "
    "papers — and what it takes to make scientific knowledge actually navigable "
    "for a research copilot. This is a 30-min talk. I'll leave time for questions. "
    "Three things I want you to leave with: (1) why search APIs aren't enough "
    "for agentic research workflows, (2) a concrete architecture that works on "
    "a single PostgreSQL box, and (3) hard-won lessons about what breaks.")

# ============================================================================
#  SLIDE 2 — The problem
# ============================================================================
s = add_blank(prs); set_bg(s)
add_title(s, "Researchers drown in literature. Agents could help — but not with today's tools.")
# big-number callouts
def stat(left, big, small, sub=None):
    add_text(s, big, left, Inches(2.2), Inches(3.5), Inches(1.3),
             size=60, bold=True, color=ACCENT_DARK, align=PP_ALIGN.CENTER)
    add_text(s, small, left, Inches(3.5), Inches(3.5), Inches(0.5),
             size=16, color=INK, align=PP_ALIGN.CENTER)
    if sub:
        add_text(s, sub, left, Inches(4.0), Inches(3.5), Inches(0.4),
                 size=12, color=INK_SOFT, italic=True, align=PP_ALIGN.CENTER)

stat(Inches(0.5), "3M+", "scientific papers published per year",
     "and doubling every ~12 years")
stat(Inches(4.9), "20–30%", "of a researcher's time", "spent on literature search")
stat(Inches(9.3), "~10", "results per search query",
     "no graph, no memory, no structure")

add_text(s, "A research copilot built on vanilla search is stuck with the same primitives as the human — scaled up.",
         Inches(0.5), Inches(5.6), Inches(12.3), Inches(0.6),
         size=18, italic=True, color=INK_MUTED, align=PP_ALIGN.CENTER)
add_footer(s, 2, TOTAL_SLIDES)
set_notes(s,
    "Scientific output is growing faster than any one researcher can keep up with. "
    "Three million papers a year across disciplines. Researchers spend 20–30% of "
    "their time on literature search according to Tenopir et al. and follow-up "
    "studies. Agents should help with this. But if you wrap Claude around ADS "
    "or Semantic Scholar's REST API, what you get is basically a faster "
    "keyword search — not actual navigation. The agent returns 10 results per "
    "query, no graph structure, no memory across queries. That's the problem I "
    "want to unpack.")

# ============================================================================
#  SLIDE 3 — What "navigable" means for an agent
# ============================================================================
s = add_blank(prs); set_bg(s)
add_title(s, "What \"navigable\" means for an agent",
          sub="Agents need different primitives than human searchers")

# Two-column comparison
col_w = Inches(5.8)
left_x = Inches(0.6)
right_x = Inches(6.9)
top = Inches(2.0)

def column(x, title, items, title_color):
    add_text(s, title, x, top, col_w, Inches(0.5),
             size=20, bold=True, color=title_color)
    add_rule(s, x, top + Inches(0.55), Inches(0.6), color=title_color)
    for i, (head, body) in enumerate(items):
        y = top + Inches(0.9 + i * 0.95)
        add_text(s, head, x, y, col_w, Inches(0.35),
                 size=15, bold=True, color=INK)
        add_text(s, body, x, y + Inches(0.35), col_w, Inches(0.55),
                 size=13, color=INK_MUTED)

column(left_x, "Human searcher",
       [("Ranked list", "Scans titles, opens a few, satisfied."),
        ("One-shot query", "Session state lives in the researcher's head."),
        ("Visual cues", "Journal, authors, year, formatting."),
        ("Stops when tired", "Implicit budget, intuition-driven.")],
       INK_SOFT)

column(right_x, "Agent",
       [("Graph topology", "Which papers bridge communities? Which are foundational?"),
        ("Session working set", "Accumulates, dedupes, tags, reasons across turns."),
        ("Structured returns", "Bibcodes, IDs, community labels, provenance."),
        ("Multi-hop moves", "Citation chains, co-citation, find_gaps, temporal.")],
       ACCENT_DARK)

add_footer(s, 3, TOTAL_SLIDES)
set_notes(s,
    "A human searcher is used to the ranked-list paradigm — you type, scan 10 "
    "results, open a couple, maybe refine. That works because the real "
    "reasoning is happening in your head. An agent has no such luxury. It needs "
    "structured primitives that compose. Graph topology tells it which papers "
    "sit at community bridges or are structurally foundational. Session state "
    "lets it accumulate and reason over a working set across dozens of tool "
    "calls. Structured returns let it track bibcodes, not free-text snippets. "
    "And it needs multi-hop moves — citation chains, finding gaps in a "
    "landscape — that are primitives the human did implicitly.")

# ============================================================================
#  SLIDE 4 — What I'll cover
# ============================================================================
s = add_blank(prs); set_bg(s)
add_title(s, "What I'll cover", sub="30 minutes · 5 sections · live demo")

sections = [
    ("1", "The problem space", "Search APIs vs. agent needs. Why MCP."),
    ("2", "Corpus completeness", "The 17.8% → 99.6% result that changes everything."),
    ("3", "Architecture", "PostgreSQL + pgvector + 13 MCP tools on a single box."),
    ("4", "Live demo", "Agent maps a research landscape with hybrid retrieval + graph moves."),
    ("5", "Hard-won lessons", "Embedding pipeline, UNLOGGED disaster, surprising eval results."),
]
for i, (num, head, body) in enumerate(sections):
    y = Inches(2.2 + i * 0.85)
    add_text(s, num, Inches(0.7), y, Inches(0.8), Inches(0.7),
             size=42, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(s, head, Inches(1.8), y + Inches(0.05), Inches(4.5), Inches(0.4),
             size=20, bold=True, color=INK)
    add_text(s, body, Inches(1.8), y + Inches(0.45), Inches(11), Inches(0.4),
             size=14, color=INK_MUTED)

add_footer(s, 4, TOTAL_SLIDES)
set_notes(s,
    "Quick roadmap. First I'll set up the problem — why search APIs are not "
    "the right interface for agentic research. Then I'll walk through the "
    "single most important empirical result we found: corpus completeness is "
    "not optional for graph analytics. Then the architecture — how this all "
    "fits on a single PostgreSQL box. A live demo walking through an actual "
    "research task. And finally the lessons I wish someone had told me before "
    "starting.")

# ============================================================================
#  SLIDE 5 — Search APIs vs. agent needs
# ============================================================================
s = add_blank(prs); set_bg(s)
add_title(s, "Search APIs return ranked lists. Agents need a knowledge structure.",
          sub="What's missing when you wrap an LLM around /search")

# Bullet grid
missing = [
    ("No graph topology",
     "No way to ask \"which papers bridge these two communities?\" or "
     "\"which papers are structurally central?\""),
    ("No session memory",
     "Every query starts cold. Working sets live in the prompt or nowhere."),
    ("No community structure",
     "Can't browse the research landscape, can't see the lens you're "
     "navigating through (citation vs. semantic vs. taxonomic)."),
    ("No multi-model signal",
     "One embedding, one ranker. Dense-only fails on specific-term queries; "
     "lexical-only fails on paraphrase."),
    ("No structured returns",
     "Free-text snippets burn context and break composition with other tools."),
]
for i, (h, b) in enumerate(missing):
    y = Inches(2.0 + i * 0.95)
    add_text(s, h, Inches(0.6), y, Inches(5), Inches(0.5),
             size=18, bold=True, color=ACCENT_DARK)
    add_text(s, b, Inches(5.8), y + Inches(0.05), Inches(7.0), Inches(0.9),
             size=14, color=INK)

add_footer(s, 5, TOTAL_SLIDES)
set_notes(s,
    "Let's be concrete about what's missing. When you wrap an LLM around ADS "
    "/search or Semantic Scholar's REST API, the agent can't ask structural "
    "questions — which papers bridge communities, which are foundational. "
    "Every query starts cold; there's no working set accumulating across "
    "turns. There's no community structure exposed, so the agent can't "
    "browse the landscape or choose a lens. Only one retrieval signal. "
    "And free-text snippets in the response burn context and break "
    "composition. All of these are primitives we'll build in the next slides.")

# ============================================================================
#  SLIDE 6 — Why MCP
# ============================================================================
s = add_blank(prs); set_bg(s)
add_title(s, "Why MCP",
          sub="Model Context Protocol is a good fit for agent-native knowledge")

reasons = [
    ("Tool-calling primitives",
     "First-class: name, schema, typed args, typed returns. The agent knows "
     "what it can do."),
    ("Session context",
     "Stateful per-connection. Working sets live on the server, not in the prompt."),
    ("Structured responses",
     "JSON objects, not HTML snippets. Compose cleanly with other tools."),
    ("Transport agnostic",
     "stdio, HTTP, WebSocket. Works with Claude Desktop, Cursor, custom agents."),
    ("Portable discipline",
     "If your research copilot talks MCP, the same knowledge layer can be "
     "hosted by ADS, arXiv, or you."),
]
for i, (h, b) in enumerate(reasons):
    y = Inches(2.0 + i * 0.95)
    add_text(s, h, Inches(0.6), y, Inches(4.3), Inches(0.5),
             size=18, bold=True, color=ACCENT_DARK)
    add_text(s, b, Inches(5.0), y + Inches(0.05), Inches(7.8), Inches(0.9),
             size=14, color=INK)

add_footer(s, 6, TOTAL_SLIDES)
set_notes(s,
    "MCP is an open protocol Anthropic introduced in late 2024, now adopted "
    "across the agent ecosystem. It gives us what we need: tool-calling "
    "primitives with schemas, stateful sessions, structured responses, "
    "transport flexibility. And importantly for scientific knowledge — it's "
    "portable. The same MCP server I built against ADS could be hosted by "
    "ADS directly, or point at arXiv, or any other corpus. The research "
    "copilot talks one protocol.")

# ============================================================================
#  SLIDE 7 — The corpus (big stats)
# ============================================================================
s = add_blank(prs); set_bg(s)
add_title(s, "The corpus",
          sub="NASA ADS — the canonical bibliographic database for astronomy")
add_image(s, ASSETS / "corpus_stats.png", Inches(0.5), Inches(2.0), width=Inches(12.3))
# secondary stats
add_text(s, "14.9M papers with full body text (46%)  ·  PostgreSQL 16 + pgvector  ·  single server (62 GB RAM)",
         Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.5),
         size=14, italic=True, color=INK_MUTED, align=PP_ALIGN.CENTER)
add_footer(s, 7, TOTAL_SLIDES)
set_notes(s,
    "Corpus at a glance. 32.4 million papers spanning 1800 to 2026. "
    "299 million citation edges. And — critically for what comes next — "
    "INDUS domain-specific embeddings for the full corpus. 14.9 million "
    "papers also have full body text ingested, which opens up chunk-level "
    "retrieval as future work. All of this runs on a single PostgreSQL 16 "
    "box with pgvector, 62 gigs of RAM. No distributed system, no separate "
    "vector database, no Elasticsearch cluster. The whole system fits in "
    "one Postgres.")

# ============================================================================
#  SLIDE 8 — Edge resolution headline
# ============================================================================
s = add_blank(prs); set_bg(s)
add_title(s, "Corpus completeness determines graph validity",
          sub="A 6-year window resolves 17.8% of citation edges. The full corpus resolves 99.6%.")
add_image(s, ASSETS / "edge_resolution.png", Inches(1.8), Inches(2.1), width=Inches(9.7))
add_text(s, "82.2% of edges dangle into the void when you index a rolling window. That makes PageRank, community detection, and co-citation analysis structurally wrong — not just incomplete.",
         Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.8),
         size=14, italic=True, color=INK_MUTED, align=PP_ALIGN.CENTER)
add_footer(s, 8, TOTAL_SLIDES)
set_notes(s,
    "OK, this is the most important empirical slide in the talk. Citation "
    "graphs are built from reference lists. Every cited paper has to exist as "
    "a node in your corpus or the edge dangles — there's no target. "
    "If you index the last 6 years of papers — which is what most commercial "
    "products do — only 17.8% of citation edges resolve. 82% of every paper's "
    "references point to older papers outside your window. "
    "The full corpus? 99.6% resolved. The remaining 0.4% are books, reports, "
    "stuff ADS doesn't track. "
    "This isn't a sampling artifact you can paper over. PageRank leaks "
    "probability through dangling edges. Leiden sees 45,000 disconnected "
    "components that shouldn't exist. Co-citation networks are sparse "
    "garbage. If you're building a research copilot and doing graph analytics "
    "on a partial corpus, you are getting structurally wrong answers.")

# ============================================================================
#  SLIDE 9 — Before/after giant component
# ============================================================================
s = add_blank(prs); set_bg(s)
add_title(s, "From 45K fragmented components to a single giant component")

# Left: 6-year window stats card
def card(left, top, w, h, fill=ACCENT_SOFT, border=ACCENT):
    r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    r.line.color.rgb = border
    r.line.width = Pt(1.2)
    r.fill.solid()
    r.fill.fore_color.rgb = fill
    return r

card(Inches(0.8), Inches(2.0), Inches(5.9), Inches(4.5),
     fill=RGBColor(0xFD, 0xEA, 0xE8), border=RED)
add_text(s, "6-year window", Inches(1.1), Inches(2.2), Inches(5.3), Inches(0.5),
         size=20, bold=True, color=RED)
add_text(s, "2021–2026, ~5M papers", Inches(1.1), Inches(2.65), Inches(5.3), Inches(0.4),
         size=13, italic=True, color=INK_MUTED)
rows_left = [("~45,000", "disconnected components"),
             ("~1.7M", "isolated nodes"),
             ("17.8%", "of edges resolved"),
             ("misleading", "PageRank, Leiden, co-citation")]
for i, (big, small) in enumerate(rows_left):
    y = Inches(3.2 + i * 0.75)
    add_text(s, big, Inches(1.1), y, Inches(2.2), Inches(0.5),
             size=22, bold=True, color=INK)
    add_text(s, small, Inches(3.4), y + Inches(0.1), Inches(3.3), Inches(0.4),
             size=13, color=INK_MUTED)

card(Inches(7.0), Inches(2.0), Inches(5.8), Inches(4.5),
     fill=RGBColor(0xE6, 0xF6, 0xEE), border=GREEN)
add_text(s, "Full corpus", Inches(7.3), Inches(2.2), Inches(5.2), Inches(0.5),
         size=20, bold=True, color=GREEN)
add_text(s, "1800–2026, 32.4M papers", Inches(7.3), Inches(2.65), Inches(5.2), Inches(0.4),
         size=13, italic=True, color=INK_MUTED)
rows_right = [("19.98M", "in the giant component (99.3% of connected)"),
              ("36", "size of the 2nd-largest component"),
              ("99.6%", "of edges resolved"),
              ("valid", "PageRank, Leiden, co-citation")]
for i, (big, small) in enumerate(rows_right):
    y = Inches(3.2 + i * 0.75)
    add_text(s, big, Inches(7.3), y, Inches(2.5), Inches(0.5),
             size=22, bold=True, color=INK)
    add_text(s, small, Inches(9.9), y + Inches(0.1), Inches(2.9), Inches(0.4),
             size=13, color=INK_MUTED)

add_text(s, "Only 1 connected component exceeds 100 nodes — the giant component itself.",
         Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
         size=13, italic=True, color=INK_SOFT, align=PP_ALIGN.CENTER)
add_footer(s, 9, TOTAL_SLIDES)
set_notes(s,
    "Same result, different angle. On the 6-year window you get ~45,000 "
    "disconnected components and 1.7M isolated nodes. That is a graph that "
    "cannot support community detection — Leiden just sees noise. "
    "On the full corpus, 99.3% of all connected papers — just under 20 million "
    "— sit in a single giant component. The second-largest component is "
    "36 nodes. Thirty-six. That's extreme bimodality, and it means community "
    "detection, PageRank, all the structural analysis — actually works. "
    "For research copilot builders: you cannot take this shortcut. If your "
    "corpus doesn't go back to the '70s at least, you are not doing graph "
    "analytics. You are doing vibes.")

# ============================================================================
#  SLIDE 10 — Takeaway for copilot builders
# ============================================================================
s = add_blank(prs); set_bg(s)
add_title(s, "Takeaway: corpus completeness is a retrieval-quality axis",
          sub="Not a \"nice to have\" — a prerequisite for agentic graph reasoning")

# three bullet cards
callouts = [
    ("Rolling-window bias",
     "Recent papers cite older papers. A rolling window systematically "
     "excludes the very papers its contents reference."),
    ("Full-corpus ingest",
     "Feasible. 32.4M papers into PostgreSQL via COPY in ~4 hours on a "
     "single machine. No excuse."),
    ("Copilot design",
     "Treat your corpus boundary as a product decision, not an operational "
     "one. It shapes what the agent can reason about."),
]
for i, (h, b) in enumerate(callouts):
    left = Inches(0.5 + i * 4.3)
    card(left, Inches(2.2), Inches(4.0), Inches(4.3),
         fill=RGBColor(0xF5, 0xF8, 0xFD), border=ACCENT)
    add_text(s, h, left + Inches(0.3), Inches(2.5), Inches(3.5), Inches(0.7),
             size=18, bold=True, color=ACCENT_DARK)
    add_text(s, b, left + Inches(0.3), Inches(3.4), Inches(3.5), Inches(2.7),
             size=14, color=INK)
add_footer(s, 10, TOTAL_SLIDES)
set_notes(s,
    "Three takeaways for anyone building a research copilot. One: rolling-window "
    "bias is systematic — recent papers cite older papers. A 2024 paper's "
    "references are mostly from 2005-2020. So indexing recent papers excludes "
    "the papers they reference. "
    "Two: the full-corpus ingest is feasible. 32 million papers into Postgres "
    "in 4 hours on a single machine with COPY. This isn't a big-data moonshot. "
    "Three: treat your corpus boundary as a product decision. It shapes what "
    "reasoning is even possible for the agent. If you only have the last 5 "
    "years, you can do topical matching but not influence, not landscape.")

# ============================================================================
#  SLIDE 11 — Architecture
# ============================================================================
s = add_blank(prs); set_bg(s)
add_title(s, "Architecture: one PostgreSQL box, three layers, one MCP server")
add_image(s, ASSETS / "architecture.png", Inches(0.4), Inches(1.7), width=Inches(12.5))
add_text(s, "Deliberate simplicity: no Elasticsearch, no separate vector DB, no orchestrator. Postgres + pgvector does it all.",
         Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.4),
         size=13, italic=True, color=INK_MUTED, align=PP_ALIGN.CENTER)
add_footer(s, 11, TOTAL_SLIDES)
set_notes(s,
    "The whole architecture. Data flows from ADS into a PostgreSQL 16 instance "
    "with pgvector. Embeddings, graph metrics, and entity extractions all "
    "live in Postgres. Above that are three horizontal layers: hybrid search, "
    "graph tools, and entity layer. Those are exposed through a 13-tool MCP "
    "server. The agent talks MCP. "
    "Deliberate simplicity. I've seen teams reach for Elasticsearch plus "
    "Pinecone plus an orchestrator, and then half the work becomes keeping "
    "those in sync. Postgres is comfortable at 32 million rows. Keep it "
    "simple until you measure that you can't.")

# ============================================================================
#  SLIDE 12 — Retrieval stack
# ============================================================================
s = add_blank(prs); set_bg(s)
add_title(s, "Retrieval stack: three signals, reciprocal rank fusion",
          sub="Dense (domain) + dense (query-doc) + lexical (BM25)")

# Three signal cards
def signal(left, title, lines, color_line):
    card(left, Inches(2.0), Inches(4.0), Inches(3.5),
         fill=RGBColor(0xF8, 0xFA, 0xFC), border=color_line)
    add_text(s, title, left + Inches(0.3), Inches(2.2), Inches(3.7), Inches(0.5),
             size=17, bold=True, color=color_line)
    for i, ln in enumerate(lines):
        add_text(s, ln, left + Inches(0.3), Inches(2.75 + i * 0.45),
                 Inches(3.6), Inches(0.4), size=13, color=INK)

signal(Inches(0.5), "INDUS",
       ["nasa-impact/nasa-smd-ibm-st-v2",
        "768d · CLS pool on title+abstract",
        "32.4M papers embedded",
        "Domain-tuned on 2.66M ADS pairs",
        "Local, open-weight, GPU only"],
       ACCENT_DARK)
signal(Inches(4.7), "text-embedding-3-large",
       ["OpenAI, 3072d → 1024d Matryoshka",
        "Query-to-document retrieval",
        "Complements INDUS on NL queries",
        "Circuit-breaker fallback",
        "Halfvec quantization"],
       ACCENT)
signal(Inches(8.9), "BM25",
       ["tsvector on title+abstract+keywords",
        "Native PostgreSQL GIN index",
        "Catches specific-term queries",
        "Cheap, interpretable",
        "Zero infra cost"],
       GREEN)

# RRF box
rrf = card(Inches(2.6), Inches(5.7), Inches(8.1), Inches(1.0),
           fill=ACCENT_SOFT, border=ACCENT_DARK)
add_text(s, "Reciprocal Rank Fusion (k=60)  →  hybrid_search",
         Inches(2.6), Inches(5.9), Inches(8.1), Inches(0.7),
         size=20, bold=True, color=ACCENT_DARK, align=PP_ALIGN.CENTER)
add_footer(s, 12, TOTAL_SLIDES)
set_notes(s,
    "The retrieval stack fuses three signals. INDUS is the NASA Science "
    "Mission Directorate's domain-tuned embedding model — trained on 2.66 "
    "million ADS title-abstract pairs. Open-weight, runs on a local GPU. "
    "We've embedded all 32 million papers with it. "
    "text-embedding-3-large is OpenAI's general-purpose model. We truncate "
    "from 3072 to 1024 dimensions using the Matryoshka property. It's "
    "designed for asymmetric query-to-document retrieval — natural-language "
    "queries hitting stored documents. INDUS alone is symmetric, so "
    "text-embedding-3-large covers a gap. "
    "BM25 — just a tsvector in Postgres — catches specific-term queries "
    "where both dense models fail. Proper names, equations, rare "
    "identifiers. "
    "All three get reciprocal-rank-fused with k=60. 49-67% error reduction "
    "over dense-only according to the literature. And we have a circuit "
    "breaker: if OpenAI's API hiccups, the stack gracefully degrades to "
    "INDUS + BM25.")

# ============================================================================
#  SLIDE 13 — Vector DB choice: pgvector + Qdrant
# ============================================================================
s = add_blank(prs); set_bg(s)
add_title(s, "pgvector + Qdrant: complementary, not competing",
          sub="Same corpus, different jobs — when each wins")

# Two-column card layout
def _vcard(left, title, sub, lines, color):
    r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.0),
                            Inches(6.0), Inches(4.7))
    r.line.color.rgb = color
    r.line.width = Pt(1.4)
    r.fill.solid()
    r.fill.fore_color.rgb = RGBColor(0xFA, 0xFB, 0xFD)
    add_text(s, title, left + Inches(0.3), Inches(2.15),
             Inches(5.5), Inches(0.5), size=20, bold=True, color=color)
    add_text(s, sub, left + Inches(0.3), Inches(2.6),
             Inches(5.5), Inches(0.35), size=12, italic=True, color=INK_SOFT)
    for i, ln in enumerate(lines):
        add_text(s, "•  " + ln, left + Inches(0.3),
                 Inches(3.1 + i * 0.55), Inches(5.6), Inches(0.5),
                 size=13, color=INK)

_vcard(Inches(0.5), "pgvector — source of truth",
       "why it won the main path",
       ["Same transaction as citation_edges, paper_metrics",
        "32M rows, 299M edges — one query plan",
        "No sync layer, no eventual consistency",
        "Halfvec + HNSW fits in one Postgres box",
        "Ops team already runs it"],
       ACCENT_DARK)

_vcard(Inches(6.8), "Qdrant — enrichment shard",
       "where it adds real value",
       ["Named multi-vectors per point (INDUS + query-doc)",
        "Payload-indexed filters at HNSW speed",
        "PQ / SQ / binary quantization + rerank",
        "Recommendation / discovery API as a primitive",
        "GPU HNSW build, per-collection snapshots"],
       ACCENT)

add_text(s,
         "Pragmatic shape: Postgres = source of truth · Qdrant = additive search/discovery layer, fed by one-way upsert",
         Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4),
         size=12, italic=True, color=INK_MUTED, align=PP_ALIGN.CENTER)
add_footer(s, 13, TOTAL_SLIDES)
set_notes(s,
    "This is a Qdrant-hosted event, so let me be direct about how I think "
    "about vector DB choice. "
    "pgvector won the main path for us because of integration, not "
    "performance. The citation graph — 299 million edges — lives in "
    "Postgres. The paper metrics — PageRank, community assignments — live "
    "in Postgres. The vectors living in the same transaction boundary "
    "means one query plan can join them cleanly. No sync layer, no "
    "eventual-consistency window. 32 million rows is comfortably inside "
    "Postgres's range. And our ops surface area is already Postgres. "
    "But there are real things Qdrant does that pgvector cannot cleanly "
    "replicate. Named multi-vectors per point — store INDUS plus "
    "text-embedding-3-large plus a chunk-level vector all on one object, "
    "query any of them. Payload-indexed filtering that keeps full HNSW "
    "speed under selective filters. Quantization options beyond halfvec. "
    "And the one I think is most interesting for agents — a first-class "
    "discovery / recommendation API: more like these, less like those. "
    "You can fake that in SQL, but it's hand-rolled. "
    "So we landed on: Postgres as source of truth; Qdrant as an additive "
    "enrichment shard, fed by one-way upsert. I'll show you the concrete "
    "capability we built on it next.")

# ============================================================================
#  SLIDE 14 — Qdrant prototype: find_similar_by_examples
# ============================================================================
s = add_blank(prs); set_bg(s)
add_title(s, "Prototype: find_similar_by_examples (14th MCP tool)",
          sub="\"More like these papers, less like those\" — optionally filtered")

# Left: signature + state
_sig_bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.95),
                              Inches(5.9), Inches(3.1))
_sig_bg.line.color.rgb = RULE; _sig_bg.line.width = Pt(1)
_sig_bg.fill.solid(); _sig_bg.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
add_text(s, "find_similar_by_examples(", Inches(0.7), Inches(2.05),
         Inches(5.5), Inches(0.35), size=13, bold=True, color=ACCENT_DARK, font="Menlo")
_sig_lines = [
    "    positive_bibcodes: list[str],",
    "    negative_bibcodes: list[str] = [],",
    "    limit: int = 10,",
    "    year_min / year_max: int | None,",
    "    doctype: list[str] | None,",
    "    community_semantic: int | None,",
    "    arxiv_class: list[str] | None,",
    ")",
]
for i, ln in enumerate(_sig_lines):
    add_text(s, ln, Inches(0.7), Inches(2.4 + i * 0.3),
             Inches(5.5), Inches(0.3), size=12, color=INK, font="Menlo")

add_text(s, "Pilot state",
         Inches(0.5), Inches(5.15), Inches(5.9), Inches(0.4),
         size=15, bold=True, color=ACCENT_DARK)
_state_rows = [
    ("400,000 points", "top-PageRank, rich payload"),
    ("INDUS 768d", "cosine similarity, HNSW"),
    ("~676 pts/s", "upsert from Postgres"),
    ("2.6 GB on disk", "one Docker container"),
]
for i, (big, small) in enumerate(_state_rows):
    y = Inches(5.55 + i * 0.35)
    add_text(s, big, Inches(0.7), y, Inches(2.2), Inches(0.3),
             size=12, bold=True, color=INK, font="Menlo")
    add_text(s, small, Inches(3.0), y, Inches(3.3), Inches(0.3),
             size=11, italic=True, color=INK_MUTED)

# Right: demo result excerpt
add_text(s, "Live demo — seed: two graphene/terahertz papers",
         Inches(6.7), Inches(2.0), Inches(6.3), Inches(0.4),
         size=13, bold=True, color=ACCENT_DARK)
_demo_lines = [
    "1. 2012NJPh...14k8001T  c=11  physics.optics       s=0.721",
    "   Comment on 'Encoding many channels on the same",
    "   frequency through radio vorticity...'",
    "2. 2010ACSNa...4.1889F  c=5   cond-mat.mes-hall   s=0.676",
    "   Thermal Conductivity of Graphene in Corbino",
    "   Membrane Geometry",
    "3. 2012NatMa..11..865V  c=15  cond-mat.mes-hall   s=0.659",
    "   Graphene FETs as room-temperature",
    "   terahertz detectors",
    "4. 2015NatNa..10..437T  c=15  cond-mat.mes-hall   s=0.658",
    "   Photovoltage in graphene on a",
    "   femtosecond timescale",
]
for i, ln in enumerate(_demo_lines):
    is_numbered = ln.lstrip().startswith(("1.", "2.", "3.", "4.", "5."))
    add_text(s, ln, Inches(6.7), Inches(2.45 + i * 0.28),
             Inches(6.3), Inches(0.3), size=11,
             bold=is_numbered,
             color=INK if is_numbered else INK_MUTED,
             font="Menlo")

add_text(s,
         "community_semantic=15 filter → restrict to graphene/cond-mat cluster.   Negative example → push results away from that direction.   All three modes, one tool call.",
         Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.5),
         size=11, italic=True, color=INK_MUTED, align=PP_ALIGN.CENTER)
add_footer(s, 14, TOTAL_SLIDES)
set_notes(s,
    "Concrete capability. This is a 14th MCP tool — only registered when "
    "QDRANT_URL is set, so the default Postgres-only deployment is "
    "unchanged. "
    "The signature: pass positive example bibcodes, optional negative "
    "example bibcodes, and optional filters — year, doctype, arxiv "
    "class, or semantic community. The handler uses Qdrant's recommend "
    "API under the hood. "
    "Pilot state: 400,000 points loaded in about 10 minutes at 676 "
    "points-per-second, pulled from the top-PageRank papers that have "
    "both a semantic-community assignment and an arXiv class. That filter "
    "combination gives us meaningful payload for demos. 2.6 GB on disk in "
    "one Docker container on the same box as Postgres. "
    "The live demo — seeding with two Nature Nanotechnology papers on "
    "graphene-based terahertz detection. The top hits are all graphene, "
    "cond-mat, optics — the discovery API found the research "
    "neighborhood cleanly. "
    "Adding a community_semantic filter restricts to the graphene "
    "cluster. Adding a negative example — like an astronomy paper — "
    "pushes results away from that direction. All three modes, one tool "
    "call. "
    "This is exactly the shape of capability that Qdrant unlocks cleanly "
    "and pgvector would have you hand-roll in SQL. Feature-flagged, "
    "additive, safe to ship alongside the Postgres source of truth.")

# ============================================================================
#  SLIDE 15 — Graph stack
# ============================================================================
s = add_blank(prs); set_bg(s)
add_title(s, "Graph intelligence: topology, communities, multi-hop moves",
          sub="Computed once on the full graph — served as columns on paper records")

col_w = Inches(4.0)
col_y = Inches(2.0)

def graph_col(left, title, lines):
    add_text(s, title, left, col_y, col_w, Inches(0.5),
             size=18, bold=True, color=ACCENT_DARK)
    add_rule(s, left, col_y + Inches(0.55), Inches(0.5))
    for i, ln in enumerate(lines):
        add_text(s, "• " + ln, left, col_y + Inches(0.9 + i * 0.55),
                 col_w, Inches(0.5), size=13, color=INK)

graph_col(Inches(0.5), "Per-paper metrics",
          ["PageRank on 32.4M nodes",
           "HITS hub / authority",
           "Three Leiden resolutions",
           "Computed weekly, ~10 min",
           "Stored as columns, not offline"])

graph_col(Inches(4.8), "Three community lenses",
          ["Citation (Leiden) — who cites whom",
           "Semantic (k-means on INDUS) — similar content",
           "Taxonomic (arXiv / UAT) — same field",
           "Disagreements are informative",
           "Agent picks the lens"])

graph_col(Inches(9.0), "Multi-hop tools",
          ["citation_chain (shortest path)",
           "citation_similarity (co-cite / biblio-couple)",
           "temporal_evolution",
           "find_gaps (bridge papers)",
           "graph_context (neighborhood)"])

add_footer(s, 15, TOTAL_SLIDES)
set_notes(s,
    "Graph intelligence layer. Three things. First, per-paper metrics: "
    "PageRank, HITS, community assignments at three resolutions — coarse, "
    "medium, fine. These are columns on the paper record in Postgres, not "
    "something computed offline and shipped around. Weekly batch, ~10 min "
    "on a tuned Postgres. "
    "Second: three orthogonal community lenses. Citation communities — "
    "Leiden algorithm on who cites whom. Semantic communities — k-means "
    "on INDUS embeddings. Taxonomic — arXiv class and UAT hierarchy. "
    "These frequently disagree, and the disagreements are informative. "
    "A paper in one citation community but a different semantic community "
    "is probably a methodological bridge. The agent picks the lens. "
    "Third: multi-hop tools. Citation chain finds the shortest citation "
    "path between two papers. Find_gaps looks for bridge papers between "
    "communities in your working set but not yet in it. These are the "
    "moves that wrap-your-LLM-around-/search just cannot make.")

# ============================================================================
#  SLIDE 16 — Tool count discipline
# ============================================================================
s = add_blank(prs); set_bg(s)
add_title(s, "13 MCP tools, down from 22",
          sub="Premortem: >15 tools degrades agent selection accuracy")
add_image(s, ASSETS / "tool_count.png", Inches(2.0), Inches(2.0), width=Inches(9.3))
# right-side list of current tools
tools = ("search · concept_search · get_paper · read_paper · citation_graph · "
         "citation_similarity · citation_chain · entity · entity_context · "
         "graph_context · find_gaps · temporal_evolution · facet_counts")
add_text(s, tools, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.7),
         size=11, italic=True, color=INK_SOFT, align=PP_ALIGN.CENTER)
add_footer(s, 16, TOTAL_SLIDES)
set_notes(s,
    "One concrete design lesson: tool count matters. We started with 22 tools "
    "across five categories. A premortem flagged that >15 tools materially "
    "degrades agent tool-selection accuracy. The A-RAG paper from Feb 2026 "
    "showed three hierarchical tools getting 94.5% on HotpotQA. "
    "We consolidated to 13. Search category went from 4 (semantic, hybrid, "
    "lexical, faceted) to 2 (search, concept_search). Session tools got "
    "folded into read-through patterns. The concrete wins: agent picks the "
    "right tool more often, context window is smaller, tool description "
    "surface area shrinks. "
    "If you're building your own MCP surface for a copilot — be ruthless. "
    "Fewer, more composable tools beat more, more specific tools.")

# ============================================================================
#  SLIDE 17 — Tool design principles
# ============================================================================
s = add_blank(prs); set_bg(s)
add_title(s, "Tool design for agents ≠ tool design for humans",
          sub="Five principles from watching Claude burn context unnecessarily")

principles = [
    ("Return compact stubs",
     "bibcode · title · first_author · year · citation_count · snippet. "
     "Full paper via read_paper only when needed."),
    ("Stable, typed interfaces",
     "JSON schema. No free-text fields. Every response parses."),
    ("Include provenance",
     "Which tool returned this? Which query? Agent can reason about its own trail."),
    ("Idempotent reads",
     "Same inputs → same outputs. Caching is trivial."),
    ("Cursor-based pagination",
     "Big result sets return a cursor, not 10,000 bibcodes."),
]
for i, (h, b) in enumerate(principles):
    y = Inches(2.0 + i * 0.95)
    add_text(s, h, Inches(0.6), y, Inches(4.3), Inches(0.5),
             size=17, bold=True, color=ACCENT_DARK)
    add_text(s, b, Inches(5.0), y + Inches(0.05), Inches(7.8), Inches(0.9),
             size=14, color=INK)
add_footer(s, 17, TOTAL_SLIDES)
set_notes(s,
    "Five principles. Compact stubs — every list response is bibcode, title, "
    "first author, year, citation count, and a short snippet. Not the full "
    "paper. The agent calls read_paper when it actually needs more. This "
    "cuts context burn by 5-10x. "
    "Stable typed interfaces — JSON schema, no free-text fields, every "
    "response parses. Otherwise the agent spends tokens parsing your HTML. "
    "Include provenance — which tool returned this, which query. Lets the "
    "agent reason about its own trail. "
    "Idempotent reads — same input, same output. Makes caching and retry "
    "trivial. "
    "Cursor-based pagination — if a query returns 10,000 bibcodes, return "
    "a cursor, not all 10,000. The agent can decide whether to keep going.")

# ============================================================================
#  SLIDE 18 — Session state architecture
# ============================================================================
s = add_blank(prs); set_bg(s)
add_title(s, "Session state: the agent's short-term memory",
          sub="In-memory working set per MCP connection — not in the prompt, not in a DB")

# ASCII-ish code block visual
code_bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(2.0),
                              Inches(6.4), Inches(4.6))
code_bg.line.color.rgb = RULE
code_bg.line.width = Pt(1)
code_bg.fill.solid()
code_bg.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)

code = (
    "WorkingSetEntry(\n"
    "    bibcode=\"2019A&A...622A...2P\",\n"
    "    source_tool=\"hybrid_search\",\n"
    "    source_context={\n"
    "        \"query\": \"exoplanet atmospheres\",\n"
    "        \"rank\": 3,\n"
    "        \"community\": \"coarse_14\"\n"
    "    },\n"
    "    tags=[\"foundational\", \"review\"],\n"
    "    added_at=\"2026-04-22T14:03:19Z\",\n"
    ")"
)
add_text(s, code, Inches(0.8), Inches(2.2), Inches(6.0), Inches(4.4),
         size=13, color=INK, font="Menlo")

# Right side: what it unlocks
right_x = Inches(7.4)
add_text(s, "What it unlocks", right_x, Inches(2.0), Inches(5.4), Inches(0.5),
         size=18, bold=True, color=ACCENT_DARK)
add_rule(s, right_x, Inches(2.55), Inches(0.5))
unlocks = [
    ("find_gaps", "Bridge papers between working-set communities not yet in the set"),
    ("get_session_summary", "What the agent knows so far, grouped by community"),
    ("Dedup and provenance", "Don't re-add. Remember how we found this."),
    ("Cross-tool reasoning", "Combine a hybrid_search result with a citation_chain hit"),
]
for i, (h, b) in enumerate(unlocks):
    y = Inches(2.9 + i * 0.9)
    add_text(s, h, right_x, y, Inches(5.4), Inches(0.35),
             size=15, bold=True, color=INK)
    add_text(s, b, right_x, y + Inches(0.35), Inches(5.4), Inches(0.5),
             size=12, color=INK_MUTED)

add_footer(s, 18, TOTAL_SLIDES)
set_notes(s,
    "Session state is the agent's short-term memory. In our MCP server it's "
    "an in-memory Python dict per stdio connection. Not persisted to a "
    "database — sessions are ephemeral by nature. Not smuggled into the "
    "prompt — that burns context. "
    "Each WorkingSetEntry has a bibcode plus provenance: which tool "
    "surfaced this, under what query, at what rank, which community. Plus "
    "agent-applied tags. "
    "What this unlocks is find_gaps — a SQL query that looks at which "
    "citation communities are represented in the working set, and finds "
    "bridge papers that connect those communities but aren't yet in the "
    "set. That's a structural move that matters for systematic reviews. "
    "Also get_session_summary grouped by community. Dedup. Cross-tool "
    "reasoning. "
    "For HTTP transport we'd persist server-side with a session ID. But "
    "for stdio — ephemeral is correct.")

# ============================================================================
#  SLIDE 19 — Three community lenses
# ============================================================================
s = add_blank(prs); set_bg(s)
add_title(s, "Three community lenses — and their informative disagreements",
          sub="Each lens has clean semantics. The agent picks.")

def lens(left, title, sub, items, color):
    card(left, Inches(2.0), Inches(4.0), Inches(4.5),
         fill=RGBColor(0xFA, 0xFB, 0xFD), border=color)
    add_text(s, title, left + Inches(0.3), Inches(2.15), Inches(3.7), Inches(0.4),
             size=18, bold=True, color=color)
    add_text(s, sub, left + Inches(0.3), Inches(2.55), Inches(3.7), Inches(0.3),
             size=11, italic=True, color=INK_SOFT)
    for i, it in enumerate(items):
        add_text(s, "•  " + it, left + Inches(0.3),
                 Inches(3.05 + i * 0.48), Inches(3.7), Inches(0.4),
                 size=13, color=INK)

lens(Inches(0.5), "Citation", "behavioral affinity",
     ["Leiden on giant component",
      "3 resolutions: ~20 · ~200 · ~2000",
      "Coverage: 62% of papers",
      "\"Papers that cite each other\"",
      "Finds invisible-college patterns"], ACCENT_DARK)
lens(Inches(4.7), "Semantic", "topical affinity",
     ["k-means on INDUS 768d",
      "3 resolutions",
      "Coverage: 99% (all INDUS-embedded)",
      "\"Papers with similar content\"",
      "Catches cross-citing subfields"], ACCENT)
lens(Inches(8.9), "Taxonomic", "institutional affinity",
     ["arXiv class + UAT concepts",
      "Hierarchical, curated",
      "Coverage: ~8% (arXiv subset)",
      "\"Papers in the same field\"",
      "Sanity-check for the other two"], GREEN)

add_footer(s, 19, TOTAL_SLIDES)
set_notes(s,
    "Three community lenses. Citation — Leiden on the giant component at "
    "three resolutions. Behavioral affinity — papers that cite each other. "
    "Semantic — k-means on INDUS embeddings. Topical affinity — similar "
    "content regardless of whether they cite each other. Taxonomic — arXiv "
    "class and UAT hierarchy. Institutional affinity — what field the "
    "paper is in by human labeling. "
    "Disagreements are the interesting cases. A paper in a citation "
    "community with a different semantic community is probably a "
    "methodological bridge — uses methods from one subfield to study "
    "problems in another. A paper with a taxonomic label disagreeing with "
    "both is often cross-disciplinary. "
    "The agent picks the lens. For a landscape-mapping task it might use "
    "taxonomic to get boundaries, then semantic to find adjacency. For "
    "a citation study, Leiden. For understanding influence, the Leiden "
    "communities plus PageRank.")

# ============================================================================
#  SLIDE 20 — Demo setup
# ============================================================================
s = add_blank(prs); set_bg(s)
add_title(s, "Live demo",
          sub="A research copilot maps the exoplanet-atmospheres landscape")

# Demo task card
dc = card(Inches(0.5), Inches(2.0), Inches(7.0), Inches(4.6),
          fill=ACCENT_SOFT, border=ACCENT)
add_text(s, "User prompt", Inches(0.8), Inches(2.2), Inches(6.4), Inches(0.4),
         size=14, bold=True, color=ACCENT_DARK)
add_text(s,
         "\"I'm preparing a grant proposal on exoplanet atmospheric "
         "characterization. Map the research landscape — what are the major "
         "subcommunities, which papers are foundational, where are the "
         "recent bridges between communities, and what am I missing?\"",
         Inches(0.8), Inches(2.7), Inches(6.4), Inches(2.2),
         size=15, italic=True, color=INK)
add_text(s, "What the copilot needs", Inches(0.8), Inches(5.0),
         Inches(6.4), Inches(0.4), size=14, bold=True, color=ACCENT_DARK)
add_text(s,
         "Topical search · community structure · PageRank within community · "
         "temporal signal · bridge-paper detection · a memory to reason over.",
         Inches(0.8), Inches(5.5), Inches(6.4), Inches(1.0),
         size=13, color=INK)

# Right side: tools invoked
add_text(s, "Tools invoked (order)", Inches(8.0), Inches(2.0),
         Inches(5.0), Inches(0.5), size=18, bold=True, color=ACCENT_DARK)
add_rule(s, Inches(8.0), Inches(2.55), Inches(0.5))
steps = [
    ("1", "search", "broad hybrid hit on \"exoplanet atmospheric characterization\""),
    ("2", "graph_context", "pull community assignments + PageRank on top hits"),
    ("3", "citation_similarity", "expand with co-citation bridges"),
    ("4", "temporal_evolution", "see which communities are growing"),
    ("5", "find_gaps", "bridge papers between working-set communities"),
    ("6", "read_paper", "deep-dive on 3 foundational papers"),
]
for i, (n, tool, body) in enumerate(steps):
    y = Inches(2.95 + i * 0.55)
    add_text(s, n, Inches(8.0), y, Inches(0.3), Inches(0.4),
             size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(s, tool, Inches(8.4), y, Inches(2.2), Inches(0.4),
             size=13, bold=True, color=INK, font="Menlo")
    add_text(s, body, Inches(10.6), y, Inches(2.5), Inches(0.4),
             size=11, color=INK_MUTED)
add_footer(s, 20, TOTAL_SLIDES)
set_notes(s,
    "OK, live demo. I'll walk through this in the next slides. "
    "The prompt is a real research copilot task: map the exoplanet "
    "atmospheric characterization landscape for a grant proposal. This "
    "needs topical search, community structure, PageRank, temporal signal, "
    "bridge-paper detection, and a memory to reason over all of it. "
    "The agent invokes six tools in this sequence. search for the broad "
    "hit. graph_context to pull community labels and PageRank. "
    "citation_similarity to expand with co-citation. temporal_evolution "
    "to see growth. find_gaps for bridge papers. And read_paper to dive "
    "deep on three foundational papers. "
    "Let me show you what this actually looks like.")

# ============================================================================
#  SLIDE 21 — Agent trace visualization
# ============================================================================
s = add_blank(prs); set_bg(s)
add_title(s, "Agent trace: tool calls visualized over time",
          sub="web/viz/agent_trace — live narration of the MCP session")
add_image(s, ASSETS / "agent_trace.png", Inches(0.8), Inches(1.9), width=Inches(11.7))
add_footer(s, 21, TOTAL_SLIDES)
set_notes(s,
    "This is the agent-trace viz from our dashboard. Each row is a tool "
    "call over time, color-coded by tool category, with the result sizes "
    "on the right. The point of this viz is to make agent behavior "
    "legible — you can see when the agent is searching broadly, when it "
    "shifts into graph moves, when it's reading papers deeply. "
    "If you're building a research copilot, building something like this "
    "for yourself is high ROI. You'll spot quickly whether the agent is "
    "using the tools you think it should, or whether it's hammering "
    "search when it should be doing graph moves. Tool-selection "
    "instrumentation is a leading indicator of copilot quality.")

# ============================================================================
#  SLIDE 22 — UMAP landscape
# ============================================================================
s = add_blank(prs); set_bg(s)
add_title(s, "UMAP of 32M INDUS embeddings — the research landscape the agent navigates",
          sub="Communities color-coded · brightness = PageRank · hover for bibcode")
add_image(s, ASSETS / "umap_browser.png", Inches(0.4), Inches(1.8), width=Inches(12.5))
add_footer(s, 22, TOTAL_SLIDES)
set_notes(s,
    "This is what the research landscape actually looks like — a UMAP "
    "projection of all 32 million INDUS embeddings. Each point is a paper. "
    "Color by coarse Leiden community, brightness by PageRank. You can see "
    "the shape of astronomy — the solar-system cluster, extragalactic, "
    "instrumentation, cosmology, each as a recognizable region. "
    "For the agent, this isn't visual — it's just community IDs and "
    "coordinates in embedding space. But for us, it's an incredibly "
    "useful QA tool. You can see when a community merge is wrong, when a "
    "semantic boundary is splitting a subfield incorrectly. "
    "The heatmap and sankey views give complementary takes — citation "
    "density between communities, flows across taxonomic boundaries.")

# ============================================================================
#  SLIDE 23 — Citation-community flow + heatmap
# ============================================================================
s = add_blank(prs); set_bg(s)
add_title(s, "Citation flow between communities",
          sub="Sankey: how citations move across Leiden communities · Heatmap: 199×199 community co-citation")
# place two screenshots side by side
add_image(s, ASSETS / "sankey.png", Inches(0.3), Inches(2.0), width=Inches(6.4))
add_image(s, ASSETS / "heatmap.png", Inches(6.9), Inches(2.0), width=Inches(6.2))
add_text(s, "These are the structural views an agent reasons over — not visually, but as graph edges and community IDs.",
         Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
         size=12, italic=True, color=INK_MUTED, align=PP_ALIGN.CENTER)
add_footer(s, 23, TOTAL_SLIDES)
set_notes(s,
    "Two more views of the same structure. On the left, a Sankey of "
    "citation flow between coarse communities — how papers in one "
    "community cite papers in another. On the right, a 199-by-199 "
    "community-level citation heatmap, medium resolution. "
    "These aren't for the agent to look at — the agent gets numbers. But "
    "they're how we as system builders validate that the community "
    "detection is meaningful and that our parallel signals agree where "
    "they should. "
    "For research copilot builders — if you can expose something like "
    "this to users, do it. It lets the human sanity-check the agent's "
    "structural claims. \"The agent says these two communities are "
    "tightly linked\" — the user can see the flow.")

# ============================================================================
#  SLIDE 24 — Retrieval eval
# ============================================================================
s = add_blank(prs); set_bg(s)
add_title(s, "Retrieval evaluation — and a surprise",
          sub="50 queries · citation-based ground truth · latest rebaseline")
add_image(s, ASSETS / "retrieval_eval.png", Inches(1.9), Inches(1.9), width=Inches(9.5))
# three quick takeaways
takes = [
    "Nomic (general-purpose) beats INDUS (domain-tuned) on this benchmark — p=0.019",
    "Hybrid adds no lift on the 10K sample: BM25 returns <12% of queries due to AND-logic",
    "All dense ≫ lexical — but each eval has methodological bias worth naming",
]
for i, t in enumerate(takes):
    add_text(s, "• " + t, Inches(0.8), Inches(6.3 + i * 0.32),
             Inches(11.8), Inches(0.3), size=12, color=INK_MUTED)
add_footer(s, 24, TOTAL_SLIDES)
set_notes(s,
    "Here's where I want to be honest with you. Our 50-query retrieval "
    "evaluation — citation-based ground truth, which means a paper is "
    "relevant if it cites or is cited by the seed paper — shows "
    "something counterintuitive. "
    "Nomic, a general-purpose embedding model with zero domain training, "
    "beats INDUS, a model fine-tuned on 2.66 million ADS title-abstract "
    "pairs. Statistically significant, p=0.019. "
    "A few hypotheses. Training-corpus diversity — Nomic's broad training "
    "may generalize better across cross-disciplinary citations. "
    "Evaluation methodology bias — citation-based ground truth captures "
    "methodological and cross-disciplinary connections that a "
    "domain-specific model optimized for within-field similarity may "
    "miss. Sample size — 50 queries with high variance. "
    "Why we still use INDUS as the primary signal: it's the only model "
    "embedded for the full 32M corpus. Nomic would cost the same "
    "compute we've already invested. And INDUS likely wins on downstream "
    "tasks beyond retrieval — community seeding, citation prediction. "
    "Lesson for copilot builders: run your own eval, pick your "
    "methodology deliberately, and name its biases.")

# ============================================================================
#  SLIDE 25 — Lesson: UNLOGGED
# ============================================================================
s = add_blank(prs); set_bg(s)
add_title(s, "Lesson 1: never use UNLOGGED tables for data you can't reconstruct",
          sub="We lost 32M SPECTER2 embeddings to a Postgres restart")

# timeline
timeline = [
    ("Day 1–3", "Bulk-embed 32M papers using UNLOGGED paper_embeddings (~15% faster writes)"),
    ("Day 4", "Postgres OOMs due to an unrelated query"),
    ("Day 4 + 1min", "systemd restarts Postgres — UNLOGGED tables are truncated at startup"),
    ("Day 4 + 2min", "32 million embeddings, ~100 GPU-hours, gone"),
    ("Day 5+", "Rebuild. Migration 023 changes paper_embeddings to LOGGED permanently."),
]
for i, (when, what) in enumerate(timeline):
    y = Inches(2.1 + i * 0.72)
    add_text(s, when, Inches(0.8), y, Inches(2.2), Inches(0.5),
             size=14, bold=True, color=RED)
    add_text(s, what, Inches(3.2), y, Inches(9.5), Inches(0.6),
             size=14, color=INK)

add_rule(s, Inches(0.8), Inches(6.0), Inches(11.7), color=RULE, height=Pt(1))
add_text(s,
         "If your pipeline builds something expensive, it must survive a "
         "process restart. UNLOGGED is a write-optimization for "
         "reconstructible data only — never persistent state.",
         Inches(0.8), Inches(6.2), Inches(11.7), Inches(1.0),
         size=13, italic=True, color=INK_MUTED)
add_footer(s, 25, TOTAL_SLIDES)
set_notes(s,
    "Hard lesson number one. When we were first embedding SPECTER2, we "
    "made paper_embeddings an UNLOGGED table for write speed — about 15% "
    "faster bulk COPY performance. We embedded 32 million papers over "
    "several days. Roughly 100 GPU-hours. "
    "Then Postgres OOM'd because of an unrelated query. systemd restarted "
    "it. And UNLOGGED tables in Postgres are truncated on crash recovery. "
    "Not persisted. Gone. "
    "32 million embeddings, 100 GPU-hours — all gone. "
    "Migration 023 fixes this permanently. paper_embeddings is LOGGED. "
    "The broader lesson: UNLOGGED is a write optimization for "
    "reconstructible data only. If you cannot cheaply re-derive it, it "
    "cannot be UNLOGGED. Every table I add now gets this decision made "
    "explicitly.")

# ============================================================================
#  SLIDE 26 — Lesson 2: pipeline optimization
# ============================================================================
s = add_blank(prs); set_bg(s)
add_title(s, "Lesson 2: the gap between naive and tuned pipelines is 10–20×",
          sub="Embedding throughput, same hardware")
add_image(s, ASSETS / "embedding_pipeline.png", Inches(2.0), Inches(1.9), width=Inches(9.3))
takes = [
    "Multiprocess: main = tokenize + GPU, writer = binary COPY",
    "UNLOGGED (for reconstructible bulk load) + drop HNSW during load, restore after",
    "Result: full 32M corpus embedded in ~24 hours on a single RTX 5090",
]
for i, t in enumerate(takes):
    add_text(s, "• " + t, Inches(0.8), Inches(6.2 + i * 0.32),
             Inches(11.8), Inches(0.3), size=13, color=INK_MUTED)
add_footer(s, 26, TOTAL_SLIDES)
set_notes(s,
    "Lesson two: the gap between a naive and a tuned embedding pipeline "
    "is an order of magnitude. We started at 32 records per second with "
    "a naive synchronous loop — tokenize, GPU inference, INSERT, repeat. "
    "Adding multiprocessing — main process does tokenization and GPU "
    "inference, a dedicated writer process does binary COPY — bumped us "
    "to 180 rec/s. "
    "Switching from individual INSERTs to buffered binary COPY — another "
    "2x to 380. "
    "Dropping the HNSW index during load and rebuilding after — final "
    "508 rec/s. "
    "Caveat on UNLOGGED from the last slide: it's fine during bulk load "
    "for reconstructible intermediate state. The lesson is about data "
    "that persists. "
    "This matters because at 32 records/sec, embedding 32 million papers "
    "takes 11 days. At 508, it's 18 hours on a single RTX 5090. That "
    "is the difference between \"feasible side project\" and \"nope\".")

# ============================================================================
#  SLIDE 27 — Lesson 3: cost estimates
# ============================================================================
s = add_blank(prs); set_bg(s)
add_title(s, "Lesson 3: LLM cost estimates are 3–5× too low",
          sub="Prompt overhead dominates output tokens for structured extraction")

# visual comparison
def cost_row(y, label, bar_width, value, color):
    add_text(s, label, Inches(0.8), y, Inches(3.2), Inches(0.5),
             size=14, bold=True, color=INK)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.0), y + Inches(0.1),
                              bar_width, Inches(0.45))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    add_text(s, value, Inches(4.0) + bar_width + Inches(0.1),
             y + Inches(0.1), Inches(3.0), Inches(0.5),
             size=14, bold=True, color=color)

add_text(s, "Entity extraction on 10K papers (Claude Haiku)",
         Inches(0.8), Inches(2.0), Inches(12), Inches(0.4),
         size=16, bold=True, color=ACCENT_DARK)
cost_row(Inches(2.6), "Naive estimate", Inches(0.7), "$10", MUTED:=RGBColor(0x9C, 0xA3, 0xAF))
cost_row(Inches(3.3), "Actual (Haiku)", Inches(3.3), "$45", RED)

add_text(s, "Full 32M corpus projection",
         Inches(0.8), Inches(4.3), Inches(12), Inches(0.4),
         size=16, bold=True, color=ACCENT_DARK)
cost_row(Inches(4.9), "Naive estimate", Inches(0.4), "$100–$3K", MUTED)
cost_row(Inches(5.6), "Premortem-corrected", Inches(2.0), "$500–$17K", RED)

add_text(s, "Why: a 500-token prompt extracting 50 tokens of structured output has 10× overhead. Caching recovers some; fine-tuning a smaller model may be necessary at full scale.",
         Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.8),
         size=12, italic=True, color=INK_MUTED)
add_footer(s, 27, TOTAL_SLIDES)
set_notes(s,
    "Lesson three, from our entity-extraction pipeline. Our initial cost "
    "estimate for running Haiku-based extraction on 10K papers was around "
    "$10. Actual was ~$45. For the full 32M corpus, naive projection was "
    "$100-$3K. Premortem corrected estimate: $500 to $17K depending on "
    "prompt design. "
    "Why the gap? Prompt overhead dominates. A 500-token system prompt "
    "plus 500 tokens of paper plus extracting 50 tokens of structured "
    "output means you're paying for 1,000 input tokens per 50 output "
    "tokens. 20x overhead. "
    "Mitigations: prompt caching (Claude supports this, cuts repeat input "
    "cost 90%). Structured output via tool-calling. And at full scale, "
    "fine-tuning a smaller open-weight model on a few thousand Haiku-"
    "labeled examples probably wins on total cost. "
    "For research copilot builders — when you're estimating LLM bills, "
    "multiply your naive estimate by 3-5. And plan for caching from day one.")

# ============================================================================
#  SLIDE 28 — What's next
# ============================================================================
s = add_blank(prs); set_bg(s)
add_title(s, "What's next",
          sub="Research copilot gets deeper as the knowledge layer gets richer")

next_items = [
    ("Chunk-and-embed full text",
     "14.9M bodies already ingested · contextual retrieval (Anthropic pattern) · "
     "35% reduction in retrieval failure projected"),
    ("Graph-RAG",
     "HippoRAG 2-style personalized PageRank over extracted concepts · "
     "7% lift over dense embeddings at 10-30x lower cost"),
    ("Better evaluation",
     "200+ query benchmark · expert-judged graded relevance · "
     "agent-task evaluation, not just retrieval"),
    ("Upstream to ADS",
     "Embeddings and graph metrics as data products · "
     "MCP endpoint alongside REST API · no re-harvest needed"),
    ("Multi-corpus",
     "Same MCP surface over arXiv, bioRxiv, APS · "
     "agent writes corpus-agnostic code"),
]
for i, (h, b) in enumerate(next_items):
    y = Inches(2.0 + i * 0.9)
    add_text(s, h, Inches(0.6), y, Inches(4.3), Inches(0.5),
             size=17, bold=True, color=ACCENT_DARK)
    add_text(s, b, Inches(5.0), y + Inches(0.05), Inches(7.8), Inches(0.9),
             size=13, color=INK)
add_footer(s, 28, TOTAL_SLIDES)
set_notes(s,
    "Where this is going. First: chunk-and-embed full text. We already have "
    "14.9 million paper bodies ingested. Applying Anthropic's contextual "
    "retrieval pattern — prepending section context to chunks — projects "
    "a 35% reduction in retrieval failure. "
    "Second: graph-RAG via HippoRAG 2 — personalized PageRank over "
    "extracted concepts. The ICML 2025 paper showed 7% lift over dense "
    "retrieval at 10-30x lower cost. Our graph primitives plus entity "
    "extraction are halfway there. "
    "Third: better evaluation. 200+ queries, graded relevance, and agent-"
    "task evaluation — not just retrieval metrics, but end-to-end copilot "
    "quality. "
    "Fourth: upstream to ADS. The long-term vision is ADS serves embeddings "
    "and graph metrics as data products — compute once, serve to everyone. "
    "And an MCP endpoint alongside their REST API so agents don't have to "
    "re-harvest the corpus. "
    "Fifth: multi-corpus. The same MCP surface should work over arXiv, "
    "bioRxiv, APS — any discipline. The agent writes corpus-agnostic code.")

# ============================================================================
#  SLIDE 29 — Checklist for research copilot builders
# ============================================================================
s = add_blank(prs); set_bg(s)
add_title(s, "A checklist for research copilot builders")

checks = [
    ("Corpus boundary",
     "Set it deliberately. Partial-corpus graph analytics are structurally wrong, not just incomplete."),
    ("Hybrid retrieval",
     "Dense (domain) + dense (query-doc) + lexical. RRF fuses. Each alone has predictable failure modes."),
    ("Graph primitives",
     "PageRank, communities, citation chains — as first-class MCP tools, not offline batch."),
    ("Session memory",
     "Working set lives on the server. Provenance tags. Not in the prompt."),
    ("Tool count discipline",
     "< 15 tools. Compact stubs. Stable schemas. Cursor pagination."),
    ("Own your eval",
     "50 queries minimum. Name your methodology bias. Validate with experts."),
    ("Instrument everything",
     "Agent-trace viz is the fastest way to learn whether your tools are working."),
]
for i, (h, b) in enumerate(checks):
    y = Inches(1.8 + i * 0.78)
    # check mark
    check = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), y + Inches(0.08),
                                Inches(0.38), Inches(0.38))
    check.line.fill.background()
    check.fill.solid()
    check.fill.fore_color.rgb = ACCENT
    add_text(s, "✓", Inches(0.5), y + Inches(0.04), Inches(0.38), Inches(0.4),
             size=16, bold=True, color=BG, align=PP_ALIGN.CENTER)
    add_text(s, h, Inches(1.05), y, Inches(4.0), Inches(0.5),
             size=16, bold=True, color=INK)
    add_text(s, b, Inches(5.3), y + Inches(0.05), Inches(7.7), Inches(0.7),
             size=13, color=INK_MUTED)
add_footer(s, 29, TOTAL_SLIDES)
set_notes(s,
    "If you take nothing else away, these seven. "
    "Set your corpus boundary deliberately — partial corpora give you "
    "structurally wrong graph analytics. "
    "Hybrid retrieval with RRF — each signal has predictable failure modes "
    "that the others cover. "
    "Graph primitives as first-class MCP tools — not an offline batch job. "
    "Session memory on the server, with provenance tags, not in the prompt. "
    "Tool count under 15. Compact stubs. Stable schemas. "
    "Own your eval. 50 queries minimum. Name your methodology bias. If you "
    "can, validate with expert judgments. "
    "And instrument everything. The agent-trace visualization I showed "
    "earlier is a 1-day build that pays dividends every time the agent "
    "behaves unexpectedly. "
    "Happy to go deeper on any of these in Q&A.")

# ============================================================================
#  SLIDE 30 — Thank you / resources
# ============================================================================
s = add_blank(prs); set_bg(s)
add_rule(s, 0, 0, Inches(0.18), height=SH, color=ACCENT)
add_text(s, "Thank you",
         Inches(0.9), Inches(1.2), Inches(12), Inches(1.4),
         size=64, bold=True, color=INK)
add_text(s, "Questions?",
         Inches(0.9), Inches(2.7), Inches(12), Inches(0.8),
         size=28, color=ACCENT_DARK)
add_rule(s, Inches(0.9), Inches(3.8), Inches(1.2))

# Contact block
add_text(s, "Stephanie Jarmak", Inches(0.9), Inches(4.1),
         Inches(12), Inches(0.5), size=22, bold=True, color=INK)
add_text(s, "Sourcegraph  ·  Harvard & Smithsonian | Center for Astrophysics",
         Inches(0.9), Inches(4.55), Inches(12), Inches(0.4),
         size=15, color=INK_MUTED)

# Resources
add_text(s, "Resources", Inches(0.9), Inches(5.3),
         Inches(12), Inches(0.4), size=16, bold=True, color=ACCENT_DARK)
resources = [
    "Repository — github.com/sjarmak/scix-agent",
    "Live MCP endpoint — mcp.sjarmak.ai/mcp/",
    "Paper (ADASS 2026, preprint) — arXiv (pending)",
    "Contact — steph.jarmak@gmail.com",
]
for i, r in enumerate(resources):
    add_text(s, "•  " + r, Inches(0.9), Inches(5.7 + i * 0.32),
             Inches(12), Inches(0.3), size=14, color=INK)

set_notes(s,
    "Thank you. Happy to take questions. "
    "The repo is open source — full code, migrations, the MCP server. "
    "The MCP server is also running publicly at mcp.sjarmak.ai — you "
    "can point Claude Desktop or Cursor at it and try it. "
    "The paper is targeting ADASS 2026, arXiv preprint dropping soon. "
    "Email is the best way to reach me.")

# ============================================================================
# SAVE
# ============================================================================
prs.save(str(OUT_PATH))
print(f"Saved {OUT_PATH} ({OUT_PATH.stat().st_size:,} bytes, {len(prs.slides)} slides)")
